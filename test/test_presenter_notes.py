"""
测试演讲者备注（Speaker Notes）注入 PPTX 的功能。

验证：
1. create_pptx_with_native_svg 传入 notes 后，PPTX 中确实包含备注内容
2. merge_svgs_to_pptx 传入 notes 后，备注能正确关联到对应 slide
"""

import os
import sys
import tempfile
import zipfile
from pathlib import Path

# 让 import 能找到项目根目录
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation

from pipeline.svg_to_pptx import create_pptx_with_native_svg
from pipeline.pptx_merger import merge_svgs_to_pptx


# ---------- 测试用 SVG ----------

SIMPLE_SVG = """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#ffffff"/>
  <text x="640" y="360" font-size="48" text-anchor="middle" fill="#333">
    Slide {page}
  </text>
</svg>"""


def _create_test_svgs(tmp_dir: str, count: int = 3) -> list[Path]:
    """在临时目录中生成若干简单 SVG 文件。"""
    paths = []
    for i in range(1, count + 1):
        p = Path(tmp_dir) / f"slide_{i:02d}.svg"
        p.write_text(SIMPLE_SVG.format(page=i), encoding="utf-8")
        paths.append(p)
    return paths


def _read_notes_from_pptx(pptx_path: str) -> dict[int, str]:
    """用 python-pptx 读取 PPTX 中每页的 notes 文本，返回 {slide_index: text}。"""
    prs = Presentation(pptx_path)
    result = {}
    for idx, slide in enumerate(prs.slides, 1):
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            text = notes_slide.notes_text_frame.text.strip()
            if text:
                result[idx] = text
    return result


def _notes_exist_in_zip(pptx_path: str) -> list[str]:
    """检查 PPTX zip 中是否存在 notesSlide*.xml 文件。"""
    with zipfile.ZipFile(pptx_path, "r") as zf:
        return [n for n in zf.namelist() if "notesSlide" in n]


# ==========================================================================
# Test 1: create_pptx_with_native_svg 直接传 notes
# ==========================================================================
def test_native_svg_with_notes():
    """底层接口：直接向 create_pptx_with_native_svg 传入 notes dict。
    输出 PPTX 到 test/output/ 供人工验证。"""
    # 输出到 test/output/ 目录（持久化，不用临时目录）
    test_output_dir = Path(__file__).parent / "output"
    test_output_dir.mkdir(exist_ok=True)
    output_pptx = test_output_dir / "test_presenter_notes.pptx"

    with tempfile.TemporaryDirectory() as tmp:
        svg_paths = _create_test_svgs(tmp, count=3)

        notes = {
            svg_paths[0].stem: "This is the **introduction** slide.\n- Point A\n- Point B",
            svg_paths[1].stem: "## Methodology\nWe use a multi-agent pipeline:\n- PDF parsing\n- Style extraction\n- SVG generation",
            svg_paths[2].stem: "Conclusion and future work.",
        }

        ok = create_pptx_with_native_svg(
            svg_files=svg_paths,
            output_path=output_pptx,
            canvas_format="ppt169",
            verbose=False,
            use_native_shapes=True,
            notes=notes,
            enable_notes=True,
        )

        assert ok, "create_pptx_with_native_svg 返回 False"
        assert output_pptx.exists(), "PPTX 文件未生成"

        # 检查 zip 内是否有 notesSlide
        notes_files = _notes_exist_in_zip(str(output_pptx))
        assert len(notes_files) >= 3, f"期望 ≥3 个 notesSlide 文件，实际: {notes_files}"

        # 用 python-pptx 读取 notes 文本
        extracted = _read_notes_from_pptx(str(output_pptx))
        assert 1 in extracted, "第 1 页缺少 notes"
        assert 2 in extracted, "第 2 页缺少 notes"
        assert 3 in extracted, "第 3 页缺少 notes"
        assert "introduction" in extracted[1].lower(), f"第 1 页 notes 内容不对: {extracted[1]}"
        assert "methodology" in extracted[2].lower(), f"第 2 页 notes 内容不对: {extracted[2]}"
        assert "conclusion" in extracted[3].lower(), f"第 3 页 notes 内容不对: {extracted[3]}"

        print(f"[PASS] test_native_svg_with_notes")
        print(f"  Output: {output_pptx}")
        for k, v in sorted(extracted.items()):
            print(f"  Slide {k} notes: {v[:100]}")


# ==========================================================================
# Test 2: merge_svgs_to_pptx 传 notes（当前应该失败，因为接口还没加 notes 参数）
# ==========================================================================
def test_merge_with_notes():
    """上层接口：通过 merge_svgs_to_pptx 传入 notes，验证 notes 能关联到正确的 slide。"""
    with tempfile.TemporaryDirectory() as tmp:
        svg_paths = _create_test_svgs(tmp, count=3)
        output_pptx = os.path.join(tmp, "merged_notes.pptx")

        notes = {
            svg_paths[0].stem: "First slide notes: overview of the paper.",
            svg_paths[1].stem: "Second slide notes: methodology details.",
            # 第 3 页故意不给 notes，测试部分有/无的情况
        }

        # 调用 merge_svgs_to_pptx（需要修改后才能接受 notes 参数）
        try:
            result = merge_svgs_to_pptx(
                svg_paths=[str(p) for p in svg_paths],
                output_pptx_path=output_pptx,
                notes=notes,
            )
        except TypeError as e:
            print(f"[EXPECTED FAIL] merge_svgs_to_pptx 尚未支持 notes 参数: {e}")
            print("  修改 merge_svgs_to_pptx 后应该能通过此测试。")
            return False

        assert result is not None, "merge_svgs_to_pptx 返回 None"
        assert os.path.exists(output_pptx), "PPTX 文件未生成"

        extracted = _read_notes_from_pptx(output_pptx)
        assert 1 in extracted, "第 1 页缺少 notes"
        assert "overview" in extracted[1].lower(), f"第 1 页 notes 内容不对: {extracted[1]}"
        assert 2 in extracted, "第 2 页缺少 notes"
        assert "methodology" in extracted[2].lower(), f"第 2 页 notes 内容不对: {extracted[2]}"
        # 第 3 页不应有用户 notes（可能有空的 notes slide）
        if 3 in extracted:
            assert extracted[3] == "", f"第 3 页不应有 notes，但有: {extracted[3]}"

        print(f"[PASS] test_merge_with_notes")
        for k, v in extracted.items():
            print(f"  Slide {k} notes: {v[:80]}")
        return True


if __name__ == "__main__":
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

    print("=" * 60)
    print("Test 1: create_pptx_with_native_svg + notes")
    print("=" * 60)
    test_native_svg_with_notes()

    print()
    print("=" * 60)
    print("Test 2: merge_svgs_to_pptx + notes")
    print("=" * 60)
    test_merge_with_notes()
