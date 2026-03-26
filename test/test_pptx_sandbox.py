from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()

# 选择空白布局
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# ======================
# 背景（用大矩形模拟渐变风格）
# ======================
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0, 0,
    prs.slide_width,
    prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(245, 230, 240)  # 淡粉色
bg.line.fill.background()

# ======================
# 标题
# ======================
title_box = slide.shapes.add_textbox(
    Inches(1), Inches(0.8),
    Inches(8), Inches(1)
)
title_tf = title_box.text_frame
title_run = title_tf.paragraphs[0].add_run()
title_run.text = "智能体技术（AI Agents）"
title_run.font.size = Pt(40)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(80, 0, 80)

# ======================
# 副标题
# ======================
subtitle_box = slide.shapes.add_textbox(
    Inches(1), Inches(1.6),
    Inches(8), Inches(0.8)
)
subtitle_tf = subtitle_box.text_frame
subtitle_run = subtitle_tf.paragraphs[0].add_run()
subtitle_run.text = "让 AI 从“工具”进化为“自主行动体”"
subtitle_run.font.size = Pt(20)
subtitle_run.font.color.rgb = RGBColor(120, 80, 120)

# ======================
# 卡片函数
# ======================
def add_card(left, top, title, content):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top,
        Inches(3.5), Inches(2)
    )
    
    # 卡片颜色
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.shadow.inherit = False
    
    # 文本
    tf = card.text_frame
    tf.clear()
    
    # 标题
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(100, 0, 100)
    
    # 内容
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(80, 80, 80)

# ======================
# 三个内容卡片
# ======================
add_card(
    Inches(1),
    Inches(2.8),
    "什么是智能体",
    "• 具备感知、决策、执行能力\n• 可自主完成任务\n• 具备环境交互能力"
)

add_card(
    Inches(5),
    Inches(2.8),
    "核心能力",
    "• 任务规划（Planning）\n• 工具调用（Tools）\n• 记忆机制（Memory）"
)

add_card(
    Inches(3),
    Inches(5),
    "应用场景",
    "• 自动办公\n• 智能编程\n• AI助手系统\n• 多智能体协作"
)

# ======================
# 底部装饰线
# ======================
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(6.8),
    prs.slide_width, Inches(0.2)
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(200, 120, 200)
line.line.fill.background()

# ======================
# 保存文件
# ======================
prs.save("AI_Agents_Presentation.pptx")

print("PPT 已生成：AI_Agents_Presentation.pptx")