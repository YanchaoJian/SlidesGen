"""
Test: 验证完整的 SVG 生成管线。

1. extract_svg_content() — 从模拟 LLM 响应中提取 SVG
2. validate_svg() — 验证 SVG 合法性
3. execute_svg() — 写入文件 + 后处理
4. merge_svgs_to_pptx() — 合并为 PPTX

Run with:
    S:/dev/miniconda3/envs/slides-gen/python.exe test/test_svg_pipeline.py
"""

import os
import sys

PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
sys.path.insert(0, PROJECT_ROOT)

from agents.composer.svg_generator import extract_svg_content
from agents.composer.svg_runner import validate_svg, execute_svg, merge_svgs_to_pptx

OUT_DIR = os.path.join(PROJECT_ROOT, "output", "test_svg_pipeline")
os.makedirs(OUT_DIR, exist_ok=True)

# ── 模拟 LLM 输出（带 markdown 代码块包裹） ──
MOCK_LLM_RESPONSE = """
Here is the SVG for slide 1:

```svg
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" width="1280" height="720">
  <defs>
    <linearGradient id="headerGrad" x1="0" y1="0" x2="1" y2="0">
      <stop offset="0%" stop-color="#003366"/>
      <stop offset="100%" stop-color="#005599"/>
    </linearGradient>
    <filter id="shadow" x="-15%" y="-15%" width="140%" height="140%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="10"/>
      <feOffset dx="0" dy="5" result="offsetBlur"/>
      <feFlood flood-color="#000000" flood-opacity="0.12" result="shadowColor"/>
      <feComposite in="shadowColor" in2="offsetBlur" operator="in" result="shadow"/>
      <feMerge>
        <feMergeNode in="shadow"/>
        <feMergeNode in="SourceGraphic"/>
      </feMerge>
    </filter>
  </defs>

  <!-- Background -->
  <rect x="0" y="0" width="1280" height="720" fill="#F5F5F5"/>

  <!-- Header bar -->
  <rect x="0" y="0" width="1280" height="100" fill="url(#headerGrad)"/>
  <text x="640" y="62" text-anchor="middle" font-family="Microsoft YaHei, Arial" font-size="36" font-weight="bold" fill="#FFFFFF">Scaled Dot-Product Attention</text>

  <!-- Content card -->
  <rect x="60" y="140" width="1160" height="480" rx="12" fill="#FFFFFF" filter="url(#shadow)"/>

  <!-- Bullet points -->
  <text x="100" y="200" font-family="Microsoft YaHei, Calibri" font-size="20" fill="#333333">
    <tspan x="100" dy="0">1. The Transformer uses scaled dot-product attention.</tspan>
    <tspan x="100" dy="36">2. Ensures stable gradients for large dimensions.</tspan>
    <tspan x="100" dy="36">3. The formula computes attention scores efficiently.</tspan>
  </text>

  <!-- Equation -->
  <text x="640" y="400" text-anchor="middle" font-family="Cambria Math, Calibri" font-size="28" fill="#003366">Attention(Q, K, V) = softmax(QK&#x1D40;/&#x221A;d&#x2096;)V</text>

  <!-- Separator -->
  <line x1="100" y1="460" x2="1180" y2="460" stroke="#003366" stroke-width="1" stroke-opacity="0.2" stroke-dasharray="4,4"/>

  <!-- Footer -->
  <text x="640" y="510" text-anchor="middle" font-family="Microsoft YaHei, Arial" font-size="14" fill="#999999">Source: Vaswani et al., "Attention Is All You Need", NeurIPS 2017</text>

  <!-- Page number -->
  <text x="1220" y="690" text-anchor="end" font-family="Arial" font-size="13" fill="#999999">04 / 10</text>
</svg>
```

This SVG creates a professional academic slide with a gradient header and content card.
"""


def test_extract():
    print("[1/4] Testing extract_svg_content()...")
    svg = extract_svg_content(MOCK_LLM_RESPONSE)
    assert svg is not None, "FAIL: extract_svg_content returned None"
    assert svg.startswith("<svg"), f"FAIL: extracted content doesn't start with <svg>: {svg[:50]}"
    assert svg.endswith("</svg>"), "FAIL: extracted content doesn't end with </svg>"
    print(f"  PASS: Extracted SVG ({len(svg)} chars)")
    return svg


def test_validate(svg):
    print("\n[2/4] Testing validate_svg()...")
    is_valid, error = validate_svg(svg)
    assert is_valid, f"FAIL: validate_svg rejected valid SVG: {error}"
    print(f"  PASS: SVG is valid")

    # Test banned feature detection
    bad_svg = '<svg xmlns="http://www.w3.org/2000/svg"><clipPath id="c"><rect/></clipPath></svg>'
    is_valid, error = validate_svg(bad_svg)
    assert not is_valid, "FAIL: validate_svg should reject SVG with clipPath"
    assert "clipPath" in error, f"FAIL: error should mention clipPath: {error}"
    print(f"  PASS: Correctly rejects banned features")


def test_execute(svg):
    print("\n[3/4] Testing execute_svg()...")
    svg_path = os.path.join(OUT_DIR, "slide_04.svg")
    success, error = execute_svg(svg, svg_path)
    assert success, f"FAIL: execute_svg failed: {error}"
    assert os.path.exists(svg_path), f"FAIL: SVG file not created: {svg_path}"

    size = os.path.getsize(svg_path)
    print(f"  PASS: SVG finalized ({size:,} bytes)")
    return svg_path


def test_merge(svg_path):
    print("\n[4/4] Testing merge_svgs_to_pptx()...")
    pptx_path = os.path.join(OUT_DIR, "test_pipeline.pptx")
    result = merge_svgs_to_pptx([svg_path], pptx_path)
    assert result is not None, "FAIL: merge_svgs_to_pptx returned None"
    assert os.path.exists(pptx_path), f"FAIL: PPTX not created: {pptx_path}"

    # Verify with python-pptx
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    shape_count = len(prs.slides[0].shapes)
    size = os.path.getsize(pptx_path)
    print(f"  PASS: PPTX created ({size:,} bytes, {slide_count} slide(s), {shape_count} shape(s))")


def main():
    print("=" * 60)
    print("SVG Pipeline Integration Test")
    print("=" * 60)

    svg = test_extract()
    test_validate(svg)
    svg_path = test_execute(svg)
    test_merge(svg_path)

    print("\n" + "=" * 60)
    print("ALL TESTS PASSED")
    print(f"Output: {OUT_DIR}")
    print("=" * 60)


if __name__ == "__main__":
    main()
