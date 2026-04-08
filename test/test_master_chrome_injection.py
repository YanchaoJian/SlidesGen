"""
Master Chrome Injection 可行性测试。

目标：验证能否把一段"chrome SVG"（页眉/页脚/logo）通过现有的
svg→DrawingML 管线转换后，注入到 PPTX slideMaster1.xml 的 spTree 中，
使得：
  1. PowerPoint 打开后每张 slide 都自动显示这段 chrome；
  2. 用户进入"视图 → 幻灯片母版"可以直接编辑 chrome 的文字/图形；
  3. 编辑保存后所有 slide 自动同步。

跑法：
    python test/test_master_chrome_injection.py

产物：
    test/output/master_chrome_test.pptx
请用 PowerPoint 打开人工验证。
"""

import os
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

# 让脚本可独立运行
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from pipeline.svg_to_pptx.drawingml_converter import convert_svg_to_slide_shapes  # noqa: E402


# ---------------------------------------------------------------------------
# 1. 测试输入：3 张已生成的 slide SVG + 一段手写的 chrome SVG
# ---------------------------------------------------------------------------

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SAMPLE_SLIDES_DIR = PROJECT_ROOT / "output" / "0408_1155_MS" / "result"
TEST_OUTPUT_DIR = PROJECT_ROOT / "test" / "output"
TEST_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 选 3 张内容页做样本（避开封面 slide_01，封面通常排版独立）
SAMPLE_SLIDE_SVGS = [
    SAMPLE_SLIDES_DIR / "slide_05" / "slide_v1.svg",
    SAMPLE_SLIDES_DIR / "slide_07" / "slide_v0.svg",
    SAMPLE_SLIDES_DIR / "slide_10" / "slide_v0.svg",
]

# 模拟 style_analyst 产出的 master chrome SVG（含一个顶部 bar + 左下校训 + 右下 logo 框）
# 这里的文字"DALIAN UNIVERSITY"和"DUT"就是用户后续在 PowerPoint master 视图直接改的内容
MASTER_CHROME_SVG = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" width="1280" height="720">
  <!-- Full-canvas background (master owns the background; slides must not draw their own) -->
  <rect x="0" y="0" width="1280" height="720" fill="#F8F9FA"/>

  <!-- Top accent bar -->
  <rect x="0" y="0" width="1280" height="6" fill="#003D7C"/>

  <!-- Bottom divider line -->
  <line x1="40" y1="660" x2="1240" y2="660" stroke="#003D7C" stroke-width="1"/>

  <!-- Footer left: institution / motto -->
  <text x="40" y="690" font-family="Arial" font-size="14" fill="#808080">
    DALIAN UNIVERSITY OF TECHNOLOGY · School of Computer Science
  </text>

  <!-- Footer right: logo placeholder box (用户在 master 里换成真 logo) -->
  <rect x="1170" y="672" width="70" height="28" fill="none" stroke="#003D7C" stroke-width="1"/>
  <text x="1205" y="692" text-anchor="middle" font-family="Arial" font-size="14" font-weight="bold" fill="#003D7C">DUT</text>

  <!-- Page number marker: will be replaced by a PPTX slide-number field after conversion -->
  <text x="1100" y="692" text-anchor="end" font-family="Arial" font-size="14" fill="#808080">PGNUM_PLACEHOLDER</text>
</svg>
"""

# 用一个不会被 XML/DrawingML 转义器拆分的纯字母 marker 标记页码位置。
# 转换完成后我们用字符串替换把它换成 <a:fld type="slidenum"> 字段。
PAGE_NUMBER_MARKER = "PGNUM_PLACEHOLDER"
# <a:fld> 必须是 <a:p> 的直接子元素，与 <a:r> 同级。所以替换粒度是整段 <a:r>...</a:r>。
# rPr 留空 (lang="en-US") 让 fld 继承段落默认样式；如需保留字体/颜色，运行时会从原 <a:r> 拷贝。
SLIDE_NUM_FIELD_TEMPLATE = (
    '<a:fld id="{{B1B2B3B4-0000-0000-0000-000000000001}}" type="slidenum">'
    '{rPr}<a:t>1</a:t></a:fld>'
)


# ---------------------------------------------------------------------------
# 工具：剥掉 slide SVG 的全画布背景 rect（以及背景图案 rect），
# 让 PPTX master 的背景透出来。
# ---------------------------------------------------------------------------

_FULL_CANVAS_RECT_RE = re.compile(
    r'<rect\b[^/>]*\bwidth\s*=\s*"1280"[^/>]*\bheight\s*=\s*"720"[^/>]*/\s*>',
    re.IGNORECASE,
)


def _strip_full_canvas_background(src: Path, dst: Path) -> None:
    """删掉 SVG 中所有 1280x720 的全画布 <rect>（通常是背景色 + 背景纹理）。

    简化策略：正则匹配 width="1280" height="720" 的自闭合 rect，全部删除。
    这能覆盖现有 slide_v*.svg 的背景写法（背景色 rect + 网格 pattern rect）。
    """
    text = src.read_text(encoding="utf-8")
    new_text, n = _FULL_CANVAS_RECT_RE.subn("", text)
    print(f"      stripped {n} full-canvas rect(s)")
    dst.write_text(new_text, encoding="utf-8")


# ---------------------------------------------------------------------------
# 2. 提取 chrome 的 spTree shapes（复用现有转换器）
# ---------------------------------------------------------------------------

def build_chrome_shapes_xml(chrome_svg_text: str, work_dir: Path) -> tuple[str, dict, list]:
    """把 chrome SVG 转成可注入 master spTree 的 shapes XML 片段。

    Returns:
        (shapes_only_xml, media_files_dict, rel_entries)
        其中 shapes_only_xml 是从转换器输出的完整 slide XML 中**只**抽出 spTree 内
        的 shape 节点（去掉 nvGrpSpPr/grpSpPr 这些容器），便于追加到 master spTree。
    """
    chrome_svg_path = work_dir / "chrome.svg"
    chrome_svg_path.write_text(chrome_svg_text, encoding="utf-8")

    full_slide_xml, media_files, rel_entries = convert_svg_to_slide_shapes(
        chrome_svg_path, slide_num=0, verbose=True
    )

    # 从完整 slide XML 中切出 <p:grpSpPr>...</p:grpSpPr> 之后到 </p:spTree> 之前的内容
    # 这就是真正的 shape 节点们
    m = re.search(
        r"</p:grpSpPr>\s*(.*?)\s*</p:spTree>",
        full_slide_xml,
        re.DOTALL,
    )
    if not m:
        raise RuntimeError("无法从转换器输出中切出 shapes 段；请检查 convert_svg_to_slide_shapes 输出格式")
    shapes_only_xml = m.group(1)

    # ── 页码 marker 替换：扫描定位法（避免正则跨元素回溯）──
    # 1. 找到 marker 字符位置
    # 2. 向左找最近的 <a:r>（中间不能出现 </a:r>，否则说明 marker 不在该 run 内）
    # 3. 向右找最近的 </a:r>
    # 4. 用 <a:fld type="slidenum"> 替换 [<a:r> ... </a:r>] 整段
    # 5. 复用原 run 的 <a:rPr>（如果存在）以保留字体/颜色
    rpr_re = re.compile(r'<a:rPr\b[^>]*(?:/>|>.*?</a:rPr>)', re.DOTALL)
    replaced = 0
    while PAGE_NUMBER_MARKER in shapes_only_xml:
        marker_pos = shapes_only_xml.find(PAGE_NUMBER_MARKER)

        # 向左找 <a:r>
        left_open = shapes_only_xml.rfind("<a:r>", 0, marker_pos)
        left_close_check = shapes_only_xml.rfind("</a:r>", 0, marker_pos)
        if left_open == -1 or (left_close_check != -1 and left_close_check > left_open):
            print(f"  ⚠️ marker 找不到包裹 <a:r>，跳出。上下文：")
            print(shapes_only_xml[max(0, marker_pos - 200):marker_pos + 200])
            break

        # 向右找 </a:r>
        right_close = shapes_only_xml.find("</a:r>", marker_pos)
        if right_close == -1:
            print("  ⚠️ marker 后找不到 </a:r>，跳出。")
            break
        right_close_end = right_close + len("</a:r>")

        run_block = shapes_only_xml[left_open:right_close_end]

        # 从 run_block 中提取 rPr（位于 <a:r> 起始之后、<a:t> 之前）
        rpr_match = rpr_re.search(run_block)
        rpr_xml = rpr_match.group(0) if rpr_match else '<a:rPr lang="en-US"/>'

        fld_xml = SLIDE_NUM_FIELD_TEMPLATE.format(rPr=rpr_xml)
        shapes_only_xml = (
            shapes_only_xml[:left_open] + fld_xml + shapes_only_xml[right_close_end:]
        )
        replaced += 1

    if replaced > 0:
        print(f"  ✓ Replaced {replaced} <a:r> run(s) containing marker with <a:fld type='slidenum'>")
    else:
        print(f"  (no '{PAGE_NUMBER_MARKER}' replacements made)")

    return shapes_only_xml, media_files, rel_entries


# ---------------------------------------------------------------------------
# 3. 注入到 slideMaster1.xml
# ---------------------------------------------------------------------------

def inject_into_master(
    extract_dir: Path,
    chrome_shapes_xml: str,
    media_files: dict,
    rel_entries: list,
) -> None:
    """把 chrome shapes 追加到 slideMaster1.xml 的 spTree，并写入 media 与 rels。"""
    master_xml_path = extract_dir / "ppt" / "slideMasters" / "slideMaster1.xml"
    master_text = master_xml_path.read_text(encoding="utf-8")

    # 把 chrome shapes 插到 master 的 </p:spTree> 之前
    if "</p:spTree>" not in master_text:
        raise RuntimeError("slideMaster1.xml 缺少 </p:spTree>，结构异常")
    new_master_text = master_text.replace(
        "</p:spTree>",
        f"\n{chrome_shapes_xml}\n</p:spTree>",
        1,
    )
    master_xml_path.write_text(new_master_text, encoding="utf-8")
    print(f"  ✓ Injected chrome shapes into {master_xml_path.relative_to(extract_dir)}")

    # 写 media（前缀避免与 slide media 撞名）
    media_dir = extract_dir / "ppt" / "media"
    media_dir.mkdir(exist_ok=True)
    for name, data in media_files.items():
        target = media_dir / f"chrome_{name}"
        target.write_bytes(data)
        print(f"  ✓ Wrote media: {target.name}")

    # 追加 master rels（如有）
    if rel_entries:
        master_rels_path = extract_dir / "ppt" / "slideMasters" / "_rels" / "slideMaster1.xml.rels"
        rels_text = master_rels_path.read_text(encoding="utf-8")
        extra = ""
        for rel in rel_entries:
            # 媒体引用要指向 chrome_ 前缀
            target = rel["target"]
            if target.startswith("../media/"):
                target = target.replace("../media/", "../media/chrome_")
            extra += (
                f'\n  <Relationship Id="{rel["id"]}" '
                f'Type="{rel["type"]}" Target="{target}"/>'
            )
        new_rels_text = rels_text.replace(
            "</Relationships>",
            f"{extra}\n</Relationships>",
            1,
        )
        master_rels_path.write_text(new_rels_text, encoding="utf-8")
        print(f"  ✓ Appended {len(rel_entries)} rel(s) to slideMaster1.xml.rels")


# ---------------------------------------------------------------------------
# 4. 完整流程
# ---------------------------------------------------------------------------

def main():
    print("=" * 70)
    print("Master Chrome Injection Test")
    print("=" * 70)

    # 校验输入
    missing = [str(p) for p in SAMPLE_SLIDE_SVGS if not p.exists()]
    if missing:
        print("❌ 缺少样本 SVG，请确认路径：")
        for m in missing:
            print(f"   - {m}")
        sys.exit(1)

    work = Path(tempfile.mkdtemp(prefix="master_chrome_test_"))
    print(f"\n[1] Work dir: {work}")

    # ── Step A: 创建 base PPTX ──
    print("\n[2] Building base PPTX with python-pptx (blank layout) ...")
    prs = Presentation()
    # 16:9 1280x720 -> EMU
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]
    for _ in SAMPLE_SLIDE_SVGS:
        prs.slides.add_slide(blank_layout)
    base_pptx = work / "base.pptx"
    prs.save(str(base_pptx))

    # ── Step B: 解压 ──
    extract_dir = work / "pptx_content"
    with zipfile.ZipFile(base_pptx, "r") as zf:
        zf.extractall(extract_dir)
    print(f"  ✓ Extracted to {extract_dir}")

    # ── Step C: 把每张 slide SVG 转为 slide XML（先剥掉全画布背景，让 master 透出来） ──
    print("\n[3] Converting slide SVGs to slide XML (stripping full-canvas backgrounds) ...")
    for i, svg_path in enumerate(SAMPLE_SLIDE_SVGS, start=1):
        print(f"  [{i}] {svg_path.name}")
        stripped_svg_path = work / f"slide_{i}_nobg.svg"
        _strip_full_canvas_background(svg_path, stripped_svg_path)
        slide_xml, media_files, rel_entries = convert_svg_to_slide_shapes(
            stripped_svg_path, slide_num=i, verbose=False
        )
        # 写 slide xml
        out = extract_dir / "ppt" / "slides" / f"slide{i}.xml"
        out.write_text(slide_xml, encoding="utf-8")

        # 写 media
        media_dir = extract_dir / "ppt" / "media"
        media_dir.mkdir(exist_ok=True)
        for name, data in media_files.items():
            (media_dir / name).write_bytes(data)

        # 写 slide rels
        rels_dir = extract_dir / "ppt" / "slides" / "_rels"
        rels_dir.mkdir(exist_ok=True)
        extra = ""
        for rel in rel_entries:
            extra += (
                f'\n  <Relationship Id="{rel["id"]}" '
                f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
            )
        rels_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\n'
            '  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
            f'{extra}\n</Relationships>'
        )
        (rels_dir / f"slide{i}.xml.rels").write_text(rels_xml, encoding="utf-8")

    # ── Step D: 把 chrome SVG 注入 master ──
    print("\n[4] Building chrome shapes from master chrome SVG ...")
    chrome_shapes_xml, chrome_media, chrome_rels = build_chrome_shapes_xml(
        MASTER_CHROME_SVG, work
    )
    print(f"  Chrome shapes XML size: {len(chrome_shapes_xml)} chars")
    print(f"  Chrome media files: {len(chrome_media)}")
    print(f"  Chrome rels: {len(chrome_rels)}")

    print("\n[5] Injecting chrome into slideMaster1.xml ...")
    inject_into_master(extract_dir, chrome_shapes_xml, chrome_media, chrome_rels)

    # ── Step E: 重新打包 ──
    output_pptx = TEST_OUTPUT_DIR / "master_chrome_test.pptx"
    if output_pptx.exists():
        output_pptx.unlink()

    print(f"\n[6] Repacking PPTX -> {output_pptx}")
    with zipfile.ZipFile(output_pptx, "w", zipfile.ZIP_DEFLATED) as zf:
        for fp in extract_dir.rglob("*"):
            if fp.is_file():
                arc = fp.relative_to(extract_dir).as_posix()
                zf.write(fp, arc)

    size = output_pptx.stat().st_size
    print(f"\n✅ Done. Output: {output_pptx} ({size:,} bytes)")
    print("\n人工验证步骤：")
    print(f"  1. 打开 {output_pptx}")
    print("  2. 翻每一页应都看到顶部蓝色 bar + 底部 'DALIAN UNIVERSITY ...' + 右下 'DUT' 框")
    print("  3. 视图 → 幻灯片母版 → 直接编辑文字/图形 → 关闭母版视图")
    print("  4. 所有 slide 应同步显示新的 chrome")

    # 不删 work 目录，便于排查
    print(f"\n(work dir 保留供调试: {work})")


if __name__ == "__main__":
    main()
