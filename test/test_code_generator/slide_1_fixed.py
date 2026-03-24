import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_slide(output_pptx_path):
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout 6 is a blank slide

    # Set slide background color to light gray
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Add a top header bar with deep blue background
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(1.5)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(26, 58, 92)
    header_bar.line.fill.background()  # Remove border

    # Add title text in the header bar
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(9),
        Inches(1)
    )
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = text_frame.add_paragraph()
    p.text = "Background: Sequence Transduction in NLP"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 245, 245)  # White text for contrast
    p.alignment = PP_ALIGN.CENTER

    # Add a text box for bullet points
    content_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.0),
        Inches(8.0),
        Inches(4.0)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Add bullet points
    bullet_points = [
        "Sequence transduction tasks, such as machine translation and language modeling, are critical in natural language processing.",
        "Historically, recurrent neural networks (RNNs) and convolutional neural networks (CNNs) have been the state-of-the-art approaches for these tasks.",
        "These models rely on sequential computation, which limits parallelization and increases training time."
    ]
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(26, 58, 92)  # Deep blue text
        p.level = 0  # Top-level bullet point

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = (
        "Introduce the audience to the importance of sequence transduction tasks in NLP "
        "and the limitations of traditional methods like RNNs and CNNs."
    )

    # Save the presentation
    prs.save(output_pptx_path)

if __name__ == "__main__":
    output_pptx_path = "/media/data_2/jyc/code/AI-Slides/SlidesGen/test/test_code_generator/slide_1.pptx"
    create_slide(output_pptx_path)