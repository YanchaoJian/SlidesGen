"""SlidesGen 统一评估框架 — LLM-as-Judge。

三个维度，每个维度一个 prompt，输出 0-5 分：
- Content（内容质量）：信息准确性、完整性、逻辑连贯性、信息密度
- Design（视觉设计）：配色对比、排版、布局对齐、视觉丰富度
- Style Transfer（风格迁移）：配色/排版/布局/装饰元素与参考风格的一致性

另附客观指标：颜色直方图相似度（HSV 空间，不依赖 LLM）。
"""

import asyncio
import json
import logging
import os
import re
import sys
from datetime import datetime
from glob import glob
from pathlib import Path

from langchain_core.messages import HumanMessage
from pptx import Presentation as PptxPresentation

# 将项目根目录加入 sys.path 以便解析 pipeline 和 utils
_PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from pipeline.pptx_imaging import pptx_to_images
from utils.llm import (
    LLMConfig,
    create_llm,
    encode_image_to_base64,
    parse_json_response,
    raise_if_fatal_llm_error,
)

logger = logging.getLogger(__name__)

# ==============================================================================
# Prompt 加载
# ==============================================================================

_PROMPT_DIR = Path(__file__).parent / "prompts"


def _load_prompt(filename: str) -> str:
    return (_PROMPT_DIR / filename).read_text(encoding="utf-8")


_CONTENT_PROMPT = _load_prompt("content.txt")
_DESIGN_PROMPT = _load_prompt("design.txt")
_STYLE_TRANSFER_PROMPT = _load_prompt("style_transfer.txt")

# ==============================================================================
# 工具函数
# ==============================================================================


def _extract_pptx_text(pptx_path: str) -> str:
    """提取 PPTX 全文，格式对齐 PPTAgent 的 Presentation.to_text()。

    输出格式：
        Slide 1 of N
        Title:{title}
        {body text}
        ----
        Slide 2 of N
        ...
    """
    prs = PptxPresentation(pptx_path)
    total = len(prs.slides)
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if (
                    slide.shapes.title
                    and shape.shape_id == slide.shapes.title.shape_id
                ):
                    title = text
                else:
                    body_parts.append(text)
        if not title and body_parts:
            title = body_parts.pop(0)
        block = f"Slide {i} of {total}\nTitle:{title}"
        if body_parts:
            block += "\n" + "\n".join(body_parts)
        blocks.append(block)
    return "\n----\n".join(blocks)


def _parse_score(response_text: str) -> dict:
    """从 LLM 响应中解析 {score, reasoning}。

    优先尝试 JSON 解析，失败则用正则提取分数。
    """
    # 尝试 JSON 解析
    parsed = parse_json_response(response_text)
    if parsed and "score" in parsed:
        return {
            "score": int(parsed["score"]),
            "reasoning": parsed.get("reasoning", ""),
        }

    # 回退：正则匹配开头的数字
    match = re.match(r"^\s*(\d+)\s*[.。]\s*(.*)", response_text, re.DOTALL)
    if match:
        return {
            "score": int(match.group(1)),
            "reasoning": match.group(2).strip(),
        }

    return {"score": 0, "reasoning": response_text}


def _compute_color_histogram_similarity(
    image_path_a: str, image_path_b: str
) -> float | None:
    """计算两张图片在 HSV 空间的颜色直方图相关系数。

    返回值范围 -1 ~ 1，越高越相似。若 OpenCV 不可用则返回 None。
    """
    try:
        import cv2
        import numpy as np
    except ImportError:
        logger.warning("OpenCV not available, skipping color histogram similarity.")
        return None

    img_a = cv2.imread(image_path_a)
    img_b = cv2.imread(image_path_b)
    if img_a is None or img_b is None:
        return None

    hsv_a = cv2.cvtColor(img_a, cv2.COLOR_BGR2HSV)
    hsv_b = cv2.cvtColor(img_b, cv2.COLOR_BGR2HSV)

    # H: 0-180, S: 0-256, V: 0-256
    h_bins, s_bins = 50, 60
    hist_size = [h_bins, s_bins]
    ranges = [0, 180, 0, 256]
    channels = [0, 1]

    hist_a = cv2.calcHist([hsv_a], channels, None, hist_size, ranges)
    hist_b = cv2.calcHist([hsv_b], channels, None, hist_size, ranges)

    cv2.normalize(hist_a, hist_a)
    cv2.normalize(hist_b, hist_b)

    return float(cv2.compareHist(hist_a, hist_b, cv2.HISTCMP_CORREL))


# ==============================================================================
# D1: 内容质量评估
# ==============================================================================


async def _eval_content(llm, presentation_text: str) -> dict:
    """评估内容质量（单次文本 LLM 调用）。"""
    logger.info("Starting D1 (Content) evaluation...")
    prompt = _CONTENT_PROMPT.replace("{{presentation_text}}", presentation_text)
    resp = await llm.ainvoke([HumanMessage(content=prompt)])
    logger.info("Finished D1 (Content) evaluation.")
    return _parse_score(resp.content.strip())


# ==============================================================================
# D2: 视觉设计质量评估
# ==============================================================================


async def _eval_design_single(vision_llm, image_path: str) -> dict:
    """评估单张 slide 截图的视觉设计质量。"""
    filename = os.path.basename(image_path)
    logger.info(f"Starting D2 (Design) evaluation for {filename}...")
    b64 = encode_image_to_base64(image_path)
    msg = HumanMessage(
        content=[
            {"type": "text", "text": _DESIGN_PROMPT},
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
        ]
    )
    resp = await vision_llm.ainvoke([msg])
    logger.info(f"Finished D2 (Design) evaluation for {filename}.")
    return _parse_score(resp.content.strip())


# ==============================================================================
# D3: 风格迁移质量评估
# ==============================================================================


async def _eval_style_transfer_single(
    vision_llm, style_image_b64: str, slide_image_path: str
) -> dict:
    """评估单张 slide 截图与参考风格图的一致性。"""
    filename = os.path.basename(slide_image_path)
    logger.info(f"Starting D3 (Style Transfer) evaluation for {filename}...")
    slide_b64 = encode_image_to_base64(slide_image_path)
    msg = HumanMessage(
        content=[
            {"type": "text", "text": _STYLE_TRANSFER_PROMPT},
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{style_image_b64}"},
            },
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{slide_b64}"},
            },
        ]
    )
    resp = await vision_llm.ainvoke([msg])
    logger.info(f"Finished D3 (Style Transfer) evaluation for {filename}.")
    return _parse_score(resp.content.strip())


# ==============================================================================
# 公开接口
# ==============================================================================


async def evaluate_pptx(
    pptx_path: str,
    llm_config: LLMConfig,
    style_image_path: str | None = None,
    output_dir: str | None = None,
    dpi: int = 200,
) -> dict:
    """对单个 PPTX 做三维统一评估。

    Args:
        pptx_path: PPTX 文件路径。
        llm_config: 视觉模型配置（需支持图片输入）。
        style_image_path: 参考风格图路径。None 则跳过 D3 风格迁移评估。
        output_dir: 评估产物输出目录。默认推导到 metrics/eval/。
        dpi: 截图 DPI。

    Returns:
        {
            "content": {"score": int, "reasoning": str},
            "design": {"score": float, "reasoning": str},
            "design_per_slide": {slide_key: {"score": int, "reasoning": str}, ...},
            "style_transfer": {"score": float, "reasoning": str} | None,
            "style_transfer_per_slide": {slide_key: {...}, ...} | None,
            "color_histogram_similarity": float | None,
        }
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    # 推导输出目录
    if output_dir is None:
        session_dir = os.path.dirname(os.path.dirname(pptx_path))
        output_dir = os.path.join(session_dir, "metrics")
    os.makedirs(output_dir, exist_ok=True)

    eval_file = os.path.join(output_dir, "eval_result.json")

    # 1. PPTX → slide 图片（只截一次）
    slide_folder = os.path.join(output_dir, "slide_images")
    if not os.path.exists(slide_folder) or not glob(
        os.path.join(slide_folder, "slide_*.jpg")
    ):
        logger.info(f"Converting PPTX to images (DPI={dpi})...")
        os.makedirs(slide_folder, exist_ok=True)
        pptx_to_images(pptx_path, slide_folder, dpi=dpi)

    slide_images = sorted(glob(os.path.join(slide_folder, "slide_*.jpg")))
    if not slide_images:
        raise RuntimeError(f"No slide images found in {slide_folder}")

    # 2. 提取全文
    logger.info("Extracting text from PPTX...")
    presentation_text = _extract_pptx_text(pptx_path)

    # 3. 创建 LLM
    logger.info("Initializing LLMs...")
    llm = create_llm(llm_config, temperature=0.0)
    vision_llm = create_llm(llm_config, temperature=0.0)

    # 4. 并行评估 D1 + D2 + D3
    # D1: 内容（单次调用）
    content_task = _eval_content(llm, presentation_text)

    # D2: 视觉设计（逐页）
    design_tasks = {
        os.path.basename(img): _eval_design_single(vision_llm, img)
        for img in slide_images
    }

    # D3: 风格迁移（逐页，需要参考图）
    style_tasks = {}
    style_image_b64 = None
    if style_image_path and os.path.exists(style_image_path):
        style_image_b64 = encode_image_to_base64(style_image_path)
        style_tasks = {
            os.path.basename(img): _eval_style_transfer_single(
                vision_llm, style_image_b64, img
            )
            for img in slide_images
        }

    # 收集所有任务并发执行
    all_keys = ["content"] + [f"design_{k}" for k in design_tasks] + [f"style_{k}" for k in style_tasks]
    all_coros = [content_task] + list(design_tasks.values()) + list(style_tasks.values())

    logger.info(
        f"Running evaluation: {len(slide_images)} slides, "
        f"D1(content) + D2(design) + {'D3(style_transfer)' if style_tasks else 'D3 skipped'}..."
    )
    all_results = await asyncio.gather(*all_coros, return_exceptions=True)

    # 5. 拆分结果
    idx = 0

    # D1
    content_result = all_results[idx]
    if isinstance(content_result, Exception):
        raise_if_fatal_llm_error(content_result)
        logger.error(f"Content evaluation failed: {content_result}")
        content_result = {"score": 0, "reasoning": f"Error: {content_result}"}
    idx += 1

    # D2
    design_per_slide = {}
    for slide_key in design_tasks:
        r = all_results[idx]
        if isinstance(r, Exception):
            raise_if_fatal_llm_error(r)
            logger.error(f"Design evaluation failed for {slide_key}: {r}")
            r = {"score": 0, "reasoning": f"Error: {r}"}
        design_per_slide[slide_key] = r
        idx += 1

    design_scores = [r["score"] for r in design_per_slide.values()]
    design_avg = round(sum(design_scores) / len(design_scores), 2) if design_scores else 0

    # D3
    style_per_slide = None
    style_avg_result = None
    color_hist_sim = None

    if style_tasks:
        style_per_slide = {}
        for slide_key in style_tasks:
            r = all_results[idx]
            if isinstance(r, Exception):
                raise_if_fatal_llm_error(r)
                logger.error(f"Style transfer evaluation failed for {slide_key}: {r}")
                r = {"score": 0, "reasoning": f"Error: {r}"}
            style_per_slide[slide_key] = r
            idx += 1

        style_scores = [r["score"] for r in style_per_slide.values()]
        style_avg = round(sum(style_scores) / len(style_scores), 2) if style_scores else 0
        style_avg_result = {"score": style_avg, "reasoning": "Average across all slides."}

        # 颜色直方图相似度（取所有页与参考图的均值）
        logger.info("Computing color histogram similarity...")
        hist_sims = []
        for img in slide_images:
            sim = _compute_color_histogram_similarity(style_image_path, img)
            if sim is not None:
                hist_sims.append(sim)
        if hist_sims:
            color_hist_sim = round(sum(hist_sims) / len(hist_sims), 4)

    # 6. 组装单次结果
    eval_result_item = {
        "evaluate_time": datetime.now().isoformat(),
        "content": content_result,
        "design": {"score": design_avg, "reasoning": "Average across all slides."},
        "style_transfer": style_avg_result,
        "color_histogram_similarity": color_hist_sim,
        "design_per_slide": design_per_slide,
        "style_transfer_per_slide": style_per_slide,
    }

    # 7. 读取旧数据、追加并持久化 (不兼容历史单次记录格式)
    history_data = {}
    if os.path.exists(eval_file):
        try:
            with open(eval_file, encoding="utf-8") as f:
                data = json.load(f)
                # 简单校验：如果根层包含 "content" 或者 "evaluation_model" 则认为是旧格式直接抛弃
                if isinstance(data, dict) and "content" not in data and "evaluation_model" not in data:
                    history_data = data
        except Exception:
            pass
            
    model_key = llm_config["model_name"] if isinstance(llm_config, dict) else getattr(llm_config, "model_name", "unknown_model")
    if model_key not in history_data:
        history_data[model_key] = []
    
    history_data[model_key].append(eval_result_item)

    with open(eval_file, "w", encoding="utf-8") as f:
        json.dump(history_data, f, indent=2, ensure_ascii=False)
    logger.info(f"Evaluation appended to {eval_file} under model '{model_key}'")

    # 8. 打印摘要
    _print_summary(eval_result_item)

    return eval_result_item


async def evaluate_pptx_batch(
    pptx_paths: list[str],
    llm_config: LLMConfig,
    style_image_paths: list[str] | None = None,
    dpi: int = 200,
) -> dict:
    """批量评估多个 PPTX，按 slide 数加权汇总。

    Args:
        pptx_paths: PPTX 文件路径列表。
        llm_config: 模型配置。
        style_image_paths: 每个 PPTX 对应的参考风格图路径列表（长度须与 pptx_paths 一致），
                           或 None（跳过风格迁移评估）。
        dpi: 截图 DPI。

    Returns:
        {"content": float, "design": float, "style_transfer": float | None,
         "color_histogram_similarity": float | None}
        分数为 0-5 量纲。
    """
    if style_image_paths and len(style_image_paths) != len(pptx_paths):
        raise ValueError("style_image_paths must have the same length as pptx_paths")

    async def _safe_eval(i: int) -> dict | None:
        style_img = style_image_paths[i] if style_image_paths else None
        try:
            return await evaluate_pptx(
                pptx_paths[i], llm_config, style_image_path=style_img, dpi=dpi
            )
        except Exception as e:
            raise_if_fatal_llm_error(e)
            logger.error(f"Failed to evaluate {pptx_paths[i]}: {e}")
            return None

    results = await asyncio.gather(*[_safe_eval(i) for i in range(len(pptx_paths))])
    valid_results = [r for r in results if r is not None]

    if not valid_results:
        return {"content": 0, "design": 0, "style_transfer": None}

    # 按 slide 数加权
    weighted = {"content": 0.0, "design": 0.0, "style_transfer": 0.0}
    hist_sims = []
    total_slides = 0
    has_style = False

    for result in valid_results:
        n = len(result.get("design_per_slide", {})) or 1
        total_slides += n
        weighted["content"] += result["content"]["score"] * n
        weighted["design"] += result["design"]["score"] * n
        if result.get("style_transfer"):
            has_style = True
            weighted["style_transfer"] += result["style_transfer"]["score"] * n
        if result.get("color_histogram_similarity") is not None:
            hist_sims.append(result["color_histogram_similarity"])

    final = {
        "content": round(weighted["content"] / total_slides, 2),
        "design": round(weighted["design"] / total_slides, 2),
        "style_transfer": round(weighted["style_transfer"] / total_slides, 2) if has_style else None,
        "color_histogram_similarity": round(sum(hist_sims) / len(hist_sims), 4) if hist_sims else None,
    }

    print("\n=== SlidesGen Evaluation Scores (0-5) ===")
    print(f"  Content:        {final['content']:.2f}")
    print(f"  Design:         {final['design']:.2f}")
    if final["style_transfer"] is not None:
        print(f"  Style Transfer: {final['style_transfer']:.2f}")
    if final["color_histogram_similarity"] is not None:
        print(f"  Color Histogram Similarity: {final['color_histogram_similarity']:.4f}")

    return final


def _print_summary(result: dict) -> None:
    """打印单次评估摘要。"""
    print("\n=== SlidesGen Evaluation Result (0-5) ===")
    print(f"  Content:        {result['content']['score']}")
    print(f"  Design:         {result['design']['score']}")
    if result.get("style_transfer"):
        print(f"  Style Transfer: {result['style_transfer']['score']}")
    if result.get("color_histogram_similarity") is not None:
        print(f"  Color Histogram: {result['color_histogram_similarity']:.4f}")


# ==============================================================================
# CLI 入口
# ==============================================================================


def main():
    """命令行入口：对单个 PPTX 执行三维评估。"""
    import argparse
    from dotenv import load_dotenv

    load_dotenv()
    logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

    parser = argparse.ArgumentParser(
        description="Evaluate a PPTX presentation (Content / Design / Style Transfer)."
    )
    parser.add_argument("--pptx_path", required=True, help="Path to the PPTX file to evaluate.")
    parser.add_argument("--style_image_path", default=None, help="Path to the reference style image (enables D3 style transfer scoring).")
    parser.add_argument("--model_name", default="gpt-4o", help="Vision model for evaluation (default: gpt-4o).")
    parser.add_argument("--output_dir", default=None, help="Directory to save evaluation results. Default: <session>/metrics/")
    parser.add_argument("--dpi", type=int, default=200, help="DPI for PPTX slide rendering (default: 200).")
    args = parser.parse_args()

    llm_config = LLMConfig(
        model_name=args.model_name,
        api_key=os.getenv("OPENAI_API_KEY", ""),
        base_url=os.getenv("OPENAI_BASE_URL", ""),
    )

    asyncio.run(
        evaluate_pptx(
            pptx_path=args.pptx_path,
            llm_config=llm_config,
            style_image_path=args.style_image_path,
            output_dir=args.output_dir,
            dpi=args.dpi,
        )
    )


if __name__ == "__main__":
    main()
