"""
PPT 风格还原脚本
从 styles.json（由 extract_ppt_style.py 生成）还原 PPTX 文件
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from typing import Dict, Any, List
import re

# EMU 单位转换常量（1英寸 = 914400 EMU）
EMU_PER_INCH = 914400
EMU_PER_POINT = 12700  # 1磅 = 12700 EMU

class PPTStyleRestorer:
    """从 styles.json 还原 PPT 文件"""
    
    def __init__(self, json_path: str):
        """
        初始化还原器
        
        Args:
            json_path: styles.json 文件路径
        """
        self.json_path = json_path
        self.style_data = self._load_json()
        
        # 初始化PPT（空白演示文稿）
        self.prs = Presentation()
        # 删除默认的空白幻灯片
        self._remove_default_slides()
        
        # 【极简兼容版】形状类型映射（仅保留最基础形状，其他全默认矩形）
        # 避免因枚举属性缺失报错
        self.shape_type_map = {
            'rect': 1,          # MSO_SHAPE.RECTANGLE 对应数值
            'roundRect': 18,    # MSO_SHAPE.ROUNDED_RECTANGLE 对应数值
            'ellipse': 9,       # MSO_SHAPE.OVAL 对应数值
            'textBox': 17,      # MSO_SHAPE.TEXT_BOX 对应数值
            'unknown': 1        # 默认矩形
        }
        
        # 对齐方式映射
        self.align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        # 垂直锚点映射（正确的枚举）
        self.vertical_anchor_map = {
            't': MSO_VERTICAL_ANCHOR.TOP,
            'b': MSO_VERTICAL_ANCHOR.BOTTOM,
            'ctr': MSO_VERTICAL_ANCHOR.MIDDLE
        }
        
        print("✅ 已加载 styles.json 并初始化 PPT")
    
    # ==================== 基础工具函数 ====================
    
    def _load_json(self) -> Dict[str, Any]:
        """加载 styles.json 文件"""
        try:
            with open(self.json_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            raise FileNotFoundError(f"未找到 JSON 文件: {self.json_path}")
        except json.JSONDecodeError:
            raise ValueError(f"JSON 文件格式错误: {self.json_path}")
    
    def _remove_default_slides(self):
        """移除PPT默认的空白幻灯片"""
        for slide in self.prs.slides:
            self.prs.slides._sldIdLst.remove(slide._sldId)
    
    def _emu_to_inches(self, emu: int) -> float:
        """EMU 转 英寸"""
        return emu / EMU_PER_INCH if emu else 0.0
    
    def _emu_to_pts(self, emu: int) -> float:
        """EMU 转 磅（Pt）"""
        return emu / EMU_PER_POINT if emu else 0.0
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """十六进制颜色转 RGBColor"""
        if not hex_color or not hex_color.startswith('#'):
            return RGBColor(0, 0, 0)  # 默认黑色
        
        hex_color = hex_color.lstrip('#')
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        except (ValueError, IndexError):
            return RGBColor(0, 0, 0)
    
    # ==================== 全局样式设置 ====================
    
    def set_global_styles(self):
        """设置全局样式（幻灯片尺寸、主题等）"""
        print("📊 开始设置全局样式...")
        
        # 1. 设置幻灯片尺寸
        self._set_slide_size()
        
        # 2. 主题颜色/字体暂不直接设置（python-pptx 对主题修改支持有限）
        global_data = self.style_data.get('global', {})
        if global_data.get('themeColors'):
            print(f"  🎨 主题颜色已读取（共{len(global_data['themeColors'])}种）")
        if global_data.get('themeFonts'):
            print(f"  🔤 主题字体已读取: {global_data['themeFonts']}")
        
        print("✅ 全局样式设置完成")
    
    def _set_slide_size(self):
        """设置幻灯片尺寸"""
        global_data = self.style_data.get('global', {})
        slide_size = global_data.get('slideSize', {})
        
        width_emu = slide_size.get('width', 9144000)  # 默认10英寸
        height_emu = slide_size.get('height', 6858000)  # 默认7.5英寸
        
        # 直接设置EMU值（最精确）
        self.prs.slide_width = Emu(width_emu)
        self.prs.slide_height = Emu(height_emu)
        
        # 打印尺寸（英寸）
        width_inch = self._emu_to_inches(width_emu)
        height_inch = self._emu_to_inches(height_emu)
        print(f"  📐 幻灯片尺寸设置为: {width_inch:.2f} × {height_inch:.2f} 英寸")
    
    # ==================== 幻灯片创建 ====================
    
    def create_slides(self):
        """创建所有幻灯片"""
        print("\n🎬 开始创建幻灯片...")
        
        slides_data = self.style_data.get('slides', [])
        if not slides_data:
            print("  ⚠️ 未找到幻灯片数据，创建空白幻灯片")
            self._create_single_slide({}, 1)
            return
        
        for idx, slide_data in enumerate(slides_data, 1):
            print(f"\n  📄 创建第 {idx} 页幻灯片...")
            self._create_single_slide(slide_data, idx)
        
        print(f"\n✅ 共创建 {len(self.prs.slides)} 页幻灯片")
    
    def _create_single_slide(self, slide_data: Dict[str, Any], slide_num: int):
        """创建单页幻灯片"""
        # 添加空白幻灯片（布局6为空白布局）
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 1. 设置幻灯片背景
        self._set_slide_background(slide, slide_data.get('background', {}))
        
        # 2. 创建所有形状
        shapes_data = slide_data.get('shapes', [])
        for shape_idx, shape_data in enumerate(shapes_data):
            try:
                self._create_shape(slide, shape_data, shape_idx)
            except Exception as e:
                print(f"    ⚠️ 创建形状 {shape_idx} 失败: {e}")
    
    def _set_slide_background(self, slide, bg_data: Dict[str, Any]):
        """设置幻灯片背景"""
        bg_type = bg_data.get('type')
        if not bg_type:
            return
        
        # 纯色背景
        if bg_type == 'solid':
            color_hex = bg_data.get('color', '#FFFFFF')
            rgb_color = self._hex_to_rgb(color_hex)
            
            # 设置背景填充
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = rgb_color
            
            # 透明度（alpha: 0-100000 → 0-1）
            alpha = bg_data.get('alpha', 100000) / 100000
            if alpha < 1.0:
                fill.fore_color.transparency = 1 - alpha
            
            print(f"    🎨 背景设置为纯色: {color_hex} (透明度: {alpha:.2f})")
        
        # 渐变背景（简化处理）
        elif bg_type == 'gradient':
            print("    ⚠️ 渐变背景暂不支持完整还原，使用默认白色背景")
    
    # ==================== 形状创建 ====================
    
    def _create_shape(self, slide, shape_data: Dict[str, Any], shape_idx: int):
        """创建单个形状（兼容所有版本）"""
        geometry = shape_data.get('geometry', {})
        pos_data = geometry.get('position', {})
        size_data = geometry.get('size', {})
        
        # 解析位置和尺寸（EMU转英寸）
        x = self._emu_to_inches(pos_data.get('x', 0))
        y = self._emu_to_inches(pos_data.get('y', 0))
        width = self._emu_to_inches(size_data.get('width', 0))
        height = self._emu_to_inches(size_data.get('height', 0))
        
        if width <= 0 or height <= 0:
            print(f"    ⚠️ 形状 {shape_idx} 尺寸无效，跳过")
            return
        
        # 解析形状类型（极简兼容：只认基础形状，其他全用矩形）
        shape_type = geometry.get('type', 'unknown')
        shape_type_val = self.shape_type_map.get(shape_type, 1)  # 1=矩形
        
        # 创建形状（直接用数值指定形状类型，避免枚举报错）
        shape = slide.shapes.add_shape(
            shape_type_val,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        shape.name = shape_data.get('name', f'Shape {shape_idx}')
        
        # 1. 设置填充样式
        self._set_shape_fill(shape, shape_data.get('fill', {}))
        
        # 2. 设置线条样式
        self._set_shape_line(shape, shape_data.get('line', {}))
        
        # 3. 设置阴影样式（兼容版）
        self._set_shape_shadow(shape, shape_data.get('shadow', {}))
        
        # 4. 设置文本内容和样式
        self._set_shape_text(shape, shape_data.get('text', {}))
    
    def _set_shape_fill(self, shape, fill_data: Dict[str, Any]):
        """设置形状填充"""
        fill_type = fill_data.get('type')
        fill = shape.fill
        
        # 无填充
        if fill_type == 'none':
            fill.background()
            return
        
        # 纯色填充
        if fill_type == 'solid':
            color_hex = fill_data.get('color', '#000000')
            rgb_color = self._hex_to_rgb(color_hex)
            
            fill.solid()
            fill.fore_color.rgb = rgb_color
            
            # 透明度
            alpha = fill_data.get('alpha', 100000) / 100000
            if alpha < 1.0:
                fill.fore_color.transparency = 1 - alpha
            return
        
        # 渐变填充（简化支持）
        if fill_type == 'gradient':
            print("    ⚠️ 渐变填充暂不支持完整还原，使用纯色替代")
            stops = fill_data.get('stops', [])
            if stops:
                color_hex = stops[0].get('color', '#FFFFFF')
                fill.solid()
                fill.fore_color.rgb = self._hex_to_rgb(color_hex)
            else:
                fill.background()
            return
        
        # 默认无填充
        fill.background()
    
    def _set_shape_line(self, shape, line_data: Dict[str, Any]):
        """设置形状线条"""
        line = shape.line
        if not line_data:
            line.fill.background()
            return
        
        # 线宽（EMU转Pt）
        width_emu = line_data.get('width', 0)
        if width_emu <= 0:
            line.fill.background()
            return
        
        # 设置线宽
        line.width = Pt(self._emu_to_pts(width_emu))
        
        # 线条颜色
        color_hex = line_data.get('color', '#000000')
        line.color.rgb = self._hex_to_rgb(color_hex)
        
        # 透明度
        alpha = line_data.get('alpha', 100000) / 100000
        if alpha < 1.0:
            line.color.transparency = 1 - alpha
        
        # 线型（虚线类型）
        dash_type = line_data.get('dashType', 'solid')
        dash_map = {
            'solid': 1,          # MSO_LINE_DASH_STYLE.SOLID 数值
            'dash': 2,           # MSO_LINE_DASH_STYLE.DASH 数值
            'dot': 3,            # MSO_LINE_DASH_STYLE.DOT 数值
            'dashDot': 4,        # MSO_LINE_DASH_STYLE.DASH_DOT 数值
            'dashDotDot': 5      # MSO_LINE_DASH_STYLE.DASH_DOT_DOT 数值
        }
        line.dash_style = dash_map.get(dash_type, 1)
    
    def _set_shape_shadow(self, shape, shadow_data: Dict[str, Any]):
        """设置形状阴影（兼容所有版本）"""
        if not shadow_data or shadow_data.get('type') not in ['outer', 'inner']:
            # 关闭阴影
            if hasattr(shape.shadow, 'inherit'):
                shape.shadow.inherit = False
            return
        
        # 启用阴影（基础兼容模式）
        shadow = shape.shadow
        shadow.inherit = True
        
        # 阴影颜色（仅基础支持）
        color_hex = shadow_data.get('color', '#000000')
        if hasattr(shadow, 'color'):
            shadow.color.rgb = self._hex_to_rgb(color_hex)
        
        # 透明度
        alpha = shadow_data.get('alpha', 100000) / 100000
        if alpha < 1.0 and hasattr(shadow, 'transparency'):
            shadow.transparency = 1 - alpha
        
        print("    ℹ️ 阴影已启用（基础模式，位置/模糊度使用默认值）")
    
    def _set_shape_text(self, shape, text_data: Dict[str, Any]):
        """设置形状文本内容和样式"""
        text_frame = shape.text_frame
        text_frame.clear()  # 清空默认文本
        
        # 设置文本框属性
        textBox_data = text_data.get('textBox', {})
        if textBox_data:
            # 内边距（EMU转英寸）
            text_frame.margin_left = Inches(self._emu_to_inches(textBox_data.get('insetLeft', 0)))
            text_frame.margin_right = Inches(self._emu_to_inches(textBox_data.get('insetRight', 0)))
            text_frame.margin_top = Inches(self._emu_to_inches(textBox_data.get('insetTop', 0)))
            text_frame.margin_bottom = Inches(self._emu_to_inches(textBox_data.get('insetBottom', 0)))
            
            # 垂直锚点（使用正确的枚举）
            anchor = textBox_data.get('anchor', 'ctr')
            text_frame.vertical_anchor = self.vertical_anchor_map.get(anchor, MSO_VERTICAL_ANCHOR.MIDDLE)
        
        # 设置段落和文本
        paragraphs = text_data.get('paragraphs', [])
        for para_data in paragraphs:
            self._add_paragraph(text_frame, para_data)
    
    def _add_paragraph(self, text_frame, para_data: Dict[str, Any]):
        """添加段落到文本框"""
        p = text_frame.add_paragraph()
        
        # 段落对齐
        align_str = para_data.get('alignment', 'left')
        p.alignment = self.align_map.get(align_str, PP_ALIGN.LEFT)
        
        # 段落左边距
        mar_l = para_data.get('marginLeft', 0)
        if mar_l > 0:
            p.left_indent = Pt(self._emu_to_pts(mar_l))
        
        # 行距
        line_spacing = para_data.get('lineSpacing')
        if line_spacing:
            if line_spacing.get('type') == 'percentage':
                p.line_spacing = line_spacing.get('value', 120) / 100
            elif line_spacing.get('type') == 'points':
                p.line_spacing = Pt(line_spacing.get('value', 12))
        
        # 添加文本段
        runs = para_data.get('runs', [])
        for run_data in runs:
            self._add_run(p, run_data)
    
    def _add_run(self, paragraph, run_data: Dict[str, Any]):
        """添加文本段到段落"""
        r = paragraph.add_run()
        r.text = run_data.get('text', '')
        
        # 字体样式
        font_data = run_data.get('font', {})
        if not font_data:
            return
        
        font = r.font
        
        # 字体大小
        font_size = font_data.get('size', 12)
        font.size = Pt(font_size)
        
        # 加粗、斜体
        font.bold = font_data.get('bold', False)
        font.italic = font_data.get('italic', False)
        
        # 下划线/删除线
        font.underline = font_data.get('underline', 'none') != 'none'
        font.strike = font_data.get('strike', 'none') != 'none'
        
        # 字体名称（优先东亚字体）
        font.name = font_data.get('ea', font_data.get('latin', '微软雅黑'))
        
        # 字体颜色
        color_hex = font_data.get('color', '#000000')
        font.color.rgb = self._hex_to_rgb(color_hex)
    
    # ==================== 保存PPT ====================
    
    def save_ppt(self, output_path: str = 'restored_ppt.pptx'):
        """保存还原后的PPT"""
        # 确保输出目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # 保存文件
        self.prs.save(output_path)
        print(f"\n✅ PPT 已保存到: {output_path}")
        return output_path

# ==================== 主程序 ====================

def main():
    """主函数"""
    import sys
    
    # 检查命令行参数
    if len(sys.argv) < 2:
        print("使用方法: python one.py <styles.json> [output.pptx]")
        print("\n示例:")
        print("  python one.py styles.json restored_ppt.pptx")
        sys.exit(1)
    
    input_json = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else 'restored_ppt.pptx'
    
    # 检查输入文件
    if not os.path.exists(input_json):
        print(f"❌ 文件不存在: {input_json}")
        sys.exit(1)
    
    print(f"📂 输入 JSON: {input_json}")
    print(f"💾 输出 PPT: {output_pptx}")
    print("=" * 60)
    
    # 创建还原器并执行
    try:
        restorer = PPTStyleRestorer(input_json)
        restorer.set_global_styles()
        restorer.create_slides()
        restorer.save_ppt(output_pptx)
        
        print("\n" + "=" * 60)
        print("✅ PPT 还原完成！")
    except Exception as e:
        print(f"\n❌ 还原失败: {e}")
        sys.exit(1)

if __name__ == '__main__':
    main()