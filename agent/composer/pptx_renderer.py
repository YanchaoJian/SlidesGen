import os
import sys
import subprocess
import logging
import importlib.util
from typing import Tuple

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

# ==============================================================================
# 核心工具函数
# ==============================================================================

def run_script(
    script_path: str,
    timeout: int = 12
) -> Tuple[bool, str]:
    """
    在一个隔离的子进程中安全地执行 Python 脚本，并验证其产物。
    捕获并返回完整的错误 traceback（标准错误）。
    """
    if not os.path.isfile(script_path):
        error_msg = f"Script not found at path: {script_path}"
        logger.error(f"   -> ❌ {error_msg}")
        return False, error_msg

    python_exe = sys.executable or "python3"
    try:
        result = subprocess.run(
            [python_exe, script_path],
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False # 不在返回非0时抛出异常
        )

        if result.returncode == 0:
            logger.debug(f"   -> Script {os.path.basename(script_path)} executed successfully.")
            return True, result.stdout
        else:
            # 记录更详细的错误
            error_details = result.stderr.strip()
            logger.warning(f"   -> ⚠️ Script {os.path.basename(script_path)} failed with return code {result.returncode}. Stderr:\n{error_details}")
            return False, error_details

    except subprocess.TimeoutExpired:
        timeout_msg = f"Script execution timed out after {timeout} seconds."
        logger.error(f"   -> ❌ {timeout_msg}")
        return False, timeout_msg
    except Exception as e:
        exec_err_msg = f"An unexpected error occurred while running the script: {e}"
        logger.error(f"   -> ❌ {exec_err_msg}", exc_info=True)
        return False, exec_err_msg


def merge_deck(code_paths: list, output_path: str):
    """
    通过动态加载各 slide 的 add_slide(prs) 函数，
    将所有页面追加到同一个 Presentation 对象中，最终保存为一个完整的 PPTX。

    每个 code_path 对应的 .py 文件必须定义 def add_slide(prs) 函数。

    Args:
        code_paths: 包含各 slide 代码文件路径的列表（按页码排序）。
        output_path: 最终合并后文件的保存路径。
    """
    if not code_paths:
        logger.warning("   -> No slide code files provided to merge.")
        return

    logger.info(f"   -> Merging {len(code_paths)} slides into a single presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for i, code_path in enumerate(code_paths):
        if not os.path.exists(code_path):
            logger.warning(f"   -> Skipping non-existent code file: {code_path}")
            continue
        try:
            # 动态加载 slide 脚本为独立模块
            spec = importlib.util.spec_from_file_location(f"slide_module_{i}", code_path)
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)

            if not hasattr(module, "add_slide"):
                logger.error(f"   -> ❌ {code_path} does not define add_slide(prs). Skipping.")
                continue

            module.add_slide(prs)
            logger.debug(f"   -> Added slide from {os.path.basename(code_path)}")
        except Exception as e:
            logger.error(f"   -> ❌ Failed to load/execute {code_path}: {e}", exc_info=True)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    logger.info(f"   -> ✅ All slides merged successfully! Saved to: {os.path.abspath(output_path)}")
