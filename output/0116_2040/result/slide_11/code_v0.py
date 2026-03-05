import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_acknowledgements_slide(prs):
    # Create a new slide
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set slide background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Draw header stripe at the top of the slide
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(255, 255, 255)
    stripe.fill.fore_color.transparency = 0.9  # 10% opacity

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_box.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = text_frame.add_paragraph()
    p.text = "ACKNOWLEDGEMENTS"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Add body text
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.625))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.line_spacing = Pt(1.2 * 16)  # Line spacing of 1.2 times font size

    # Add bullet points
    bullet_points = [
        "Thank funding agencies, collaborators, advisors, and institutions.",
        "Include logos of funding sources if available."
    ]
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.space_after = Pt(12)  # Space after each bullet point
        p.level = 0  # Bullet point level

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Appropriately acknowledge all contributions and support."

if __name__ == "__main__":
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Create the acknowledgements slide
    create_acknowledgements_slide(prs)

    # Define output path
    output_pptx_path = 'output/0116_2040/result/slide_11/slide.pptx'
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)

    # Save the presentation
    prs.save(output_pptx_path)