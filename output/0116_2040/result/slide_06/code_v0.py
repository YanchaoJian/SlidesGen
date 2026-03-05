import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_slide(prs):
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Add header stripe
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.fore_color.transparency = 0.9  # 0.1 opacity

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.word_wrap = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title = title_frame.add_paragraph()
    title.text = "EXPERIMENTAL SETUP"
    title.font.size = Pt(32)
    title.font.bold = True
    title.font.color.rgb = RGBColor(255, 255, 255)
    title.font.name = "Arial"

    # Add bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.625))
    content_frame = content_box.text_frame
    content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    content_frame.word_wrap = True
    content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    bullet_points = [
        "Experiments conducted on WMT 2014 English-to-German and English-to-French translation tasks.",
        "Evaluation metrics include BLEU scores.",
        "Baseline methods used for comparison include state-of-the-art models."
    ]

    for point in bullet_points:
        p = content_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.font.name = "Arial"
        p.space_after = Pt(14)  # Line spacing

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Introduce the experimental setup and evaluation metrics used in the study."

if __name__ == "__main__":
    output_pptx_path = 'output/0116_2040/result/slide_06/slide.pptx'
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    create_slide(prs)
    prs.save(output_pptx_path)