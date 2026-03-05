import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_pptx_path):
    # Create a presentation object
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add a slide with a blank layout
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)
    
    # Add a header stripe at the top of the slide
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # 10% opacity
    
    # Add the title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = "Solution Overview: The Transformer Model"
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.size = Pt(32)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.name = 'Arial'
    text_frame.paragraphs[0].text = text_frame.paragraphs[0].text.upper()
    
    # Add the bullet points
    body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(3.0))
    text_frame = body_shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = text_frame.add_paragraph()
    p.text = "The Transformer model introduces a novel architecture based solely on attention mechanisms."
    p.space_after = Pt(14.4)  # Line spacing 1.2
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(163, 185, 204)
    p.font.name = 'Arial'
    
    p = text_frame.add_paragraph()
    p.text = "It eliminates the need for recurrence and convolution, allowing for greater parallelization."
    p.space_after = Pt(14.4)
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(163, 185, 204)
    p.font.name = 'Arial'
    
    p = text_frame.add_paragraph()
    p.text = "The model sets new benchmarks in machine translation tasks."
    p.space_after = Pt(14.4)
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(163, 185, 204)
    p.font.name = 'Arial'
    
    # Add the image
    image_path = "path/to/transformer_model_architecture.png"  # Update with the correct path
    slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4.0), height=Inches(3.0))
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Provide an overview of the Transformer model and its advantages."
    
    # Save the presentation
    prs.save(output_pptx_path)

if __name__ == "__main__":
    output_pptx_path = "output/0116_2040/result/slide_03/slide.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    create_presentation(output_pptx_path)