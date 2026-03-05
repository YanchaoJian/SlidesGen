import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_pptx_path):
    # Create a presentation object with a 16:9 slide layout
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set slide background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Add a header stripe at the top of the slide
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # 10% opacity

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = "Background: Sequence Transduction Models"
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.size = Pt(32)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.name = 'Arial'
    text_frame.paragraphs[0].text = text_frame.paragraphs[0].text.upper()

    # Add bullet points
    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.625))
    text_frame = content_shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    bullet_points = [
        "Sequence transduction models are essential for tasks like language modeling and machine translation.",
        "Traditionally, these models rely on recurrent neural networks (RNNs) and convolutional neural networks (CNNs).",
        "Attention mechanisms have emerged as a powerful tool to model dependencies in sequences."
    ]

    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.font.name = 'Arial'
        p.space_after = Pt(14.4)  # Line spacing of 1.2 * 12pt

    # Add presenter notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Introduce the importance of sequence transduction models and the role of attention mechanisms."

    # Save the presentation
    prs.save(output_pptx_path)

if __name__ == "__main__":
    output_pptx_path = "output/0116_2040/result/slide_01/slide.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    create_presentation(output_pptx_path)