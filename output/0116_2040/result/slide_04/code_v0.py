import os
import io
import tempfile
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def render_latex_to_image(latex_str, dpi=300):
    """
    Renders a single-line LaTeX string to an image file using matplotlib.
    No \begin{{}} blocks allowed.
    """
    clean_latex = latex_str.strip().strip('$')
    wrapped_latex = f"${clean_latex}$"
    
    fig, ax = plt.subplots(figsize=(0.1, 0.1))
    ax.axis('off')
    
    try:
        ax.text(0.5, 0.5, wrapped_latex, size=20, ha='center', va='center')
        
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            plt.savefig(tmp.name, format='png', bbox_inches='tight', pad_inches=0.1, dpi=dpi)
            return tmp.name
    except Exception as e:
        print(f"LaTeX render error: {e}")
        plt.close(fig)
        return None
    finally:
        plt.close(fig)

def create_presentation():
    # Create a presentation object
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set slide dimensions to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Add header stripe
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # 10% opacity

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Technical Approach: Multi-Head Self-Attention"
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].font.name = 'Arial'
    title_frame.paragraphs[0].text = title_frame.paragraphs[0].text.upper()

    # Add bullet points
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(3))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True
    bullet_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_frame.paragraphs[0].font.size = Pt(16)
    bullet_frame.paragraphs[0].font.color.rgb = RGBColor(163, 185, 204)
    bullet_frame.paragraphs[0].font.name = 'Arial'
    bullet_frame.paragraphs[0].line_spacing = Pt(19.2)  # 1.2 line spacing

    # Add bullet points text
    bullet_points = [
        "Point 1: Overview of the mechanism.",
        "Point 2: Key advantages.",
        "Point 3: Implementation details."
    ]
    for point in bullet_points:
        p = bullet_frame.add_paragraph()
        p.text = point
        p.level = 0

    # Add image
    image_path = "output/0116_2040/images/_page_3_Figure_1.jpeg"
    slide.shapes.add_picture(image_path, Inches(6), Inches(1.5), height=Inches(3))

    # Add equation
    equation = "Attention(Q, K, V) = \\operatorname{softmax}(\\frac{QK^T}{\\sqrt{d_k}})V"
    equation_image_path = render_latex_to_image(equation)
    if equation_image_path:
        slide.shapes.add_picture(equation_image_path, Inches(3), Inches(4.5), height=Inches(1))

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Explain the technical details of the attention mechanism used in the Transformer."

    return prs

if __name__ == "__main__":
    output_pptx_path = "output/0116_2040/result/slide_04/slide.pptx"
    prs = create_presentation()
    prs.save(output_pptx_path)