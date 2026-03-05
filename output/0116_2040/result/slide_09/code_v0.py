import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_slide(output_pptx_path):
    # Create a presentation object
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add a slide with a blank layout
    slide_layout = prs.slide_layouts[5]  # 5 is typically a blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)
    
    # Add a header stripe at the top
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # Set opacity to 0.1 (transparency to 0.9)
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "CONCLUSION & IMPACT"
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].font.name = 'Arial'
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Add bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.625))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    bullet_points = [
        "Summarize the impact of the Transformer model.",
        "Discuss future research directions.",
        "Highlight key findings and implications."
    ]
    
    for point in bullet_points:
        p = content_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.font.name = 'Arial'
        p.space_after = Pt(14.4)  # Line spacing of 1.2 * 12pt (default line height)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Summarize the impact of the Transformer model and future research directions."

    # Save the presentation
    prs.save(output_pptx_path)

if __name__ == "__main__":
    output_pptx_path = "output/0116_2040/result/slide_09/slide.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    create_slide(output_pptx_path)