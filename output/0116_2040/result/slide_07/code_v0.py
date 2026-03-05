import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_slide(prs):
    # Create a new slide
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Draw a header stripe at the top
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # Set opacity to 0.1

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "Key Results: BLEU Scores and Training Costs"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT

    # Add bullet points
    bullet_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    text_frame = bullet_shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_points = [
        "The Transformer model achieves superior BLEU scores compared to previous models.",
        "Significant reduction in training costs observed.",
        "Results demonstrate the efficiency and effectiveness of the Transformer model."
    ]
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.font.name = 'Arial'
        p.space_after = Pt(14)  # Line spacing

    # Add table
    rows, cols = 3, 3  # Example table size, adjust as needed
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.0), Inches(9), Inches(2)).table
    for row in range(rows):
        for col in range(cols):
            cell = table_shape.cell(row, col)
            cell.text = f"Row {row+1}, Col {col+1}"
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(163, 185, 204)
            cell.text_frame.paragraphs[0].font.name = 'Arial'

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Highlight the superior performance and reduced training costs of the Transformer model."

if __name__ == "__main__":
    output_pptx_path = 'output/0116_2040/result/slide_07/slide.pptx'
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    create_slide(prs)
    prs.save(output_pptx_path)