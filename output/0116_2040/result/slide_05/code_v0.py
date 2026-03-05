import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import tempfile

def render_latex_to_image(latex_str, dpi=300):
    """
    Renders a single-line LaTeX string to an image file using matplotlib.
    No \begin{{}} blocks allowed.
    """
    clean_latex = latex_str.strip().strip('$')
    wrapped_latex = f"${clean_latex}$"
    
    fig, ax = plt.subplots(figsize=(0.1, 0.1))
    ax.axis('off')
    
    try:
        ax.text(0.5, 0.5, wrapped_latex, size=20, ha='center', va='center')
        
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            plt.savefig(tmp.name, format='png', bbox_inches='tight', pad_inches=0.1, dpi=dpi)
            return tmp.name
    except Exception as e:
        print(f"LaTeX render error: {e}")
        plt.close(fig)
        return None
    finally:
        plt.close(fig)

def create_presentation(output_pptx_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set slide background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Add header stripe
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # 10% opacity

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "POSITIVE ENCODING IN TRANSFORMERS"
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.name = 'Arial'
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add bullet points
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True
    bullet_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_frame.paragraphs[0].space_after = Pt(0)
    bullet_frame.paragraphs[0].space_before = Pt(0)
    bullet_frame.paragraphs[0].line_spacing = 1.2

    bullet_points = [
        "Positional encoding helps the model understand the order of tokens in a sequence.",
        "Sine and cosine functions are used for even and odd indices, respectively.",
        "This encoding is crucial for maintaining sequence information without recurrence."
    ]

    for point in bullet_points:
        p = bullet_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = 1.2

    # Add LaTeX formula
    latex_formula = "PE_{(pos,2i)} = \\sin(pos/10000^{2i/d_{\\text{model}}})"
    formula_image_path = render_latex_to_image(latex_formula)

    if formula_image_path:
        slide.shapes.add_picture(formula_image_path, Inches(0.5), Inches(3.7), width=Inches(9))
        os.remove(formula_image_path)

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Describe how positional encoding works and its importance in the Transformer model."

    prs.save(output_pptx_path)

if __name__ == "__main__":
    output_pptx_path = "output/0116_2040/result/slide_05/slide.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    create_presentation(output_pptx_path)