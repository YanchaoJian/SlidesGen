import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_pptx_path):
    # Create a presentation object
    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Add a slide with a blank layout
    slide_layout = prs.slide_layouts[5]  # 5 is typically a blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 31, 68)

    # Add a header stripe at the top
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(255, 255, 255)
    header_fill.transparency = 0.9  # Set opacity to 10%

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_box.text_frame
    text_frame.text = "Analysis: Generalization to Other Tasks"
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.size = Pt(32)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.name = 'Arial'
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add bullet points
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True
    bullet_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_frame.paragraphs[0].font.size = Pt(16)
    bullet_frame.paragraphs[0].font.color.rgb = RGBColor(163, 185, 204)
    bullet_frame.paragraphs[0].font.name = 'Arial'
    bullet_frame.paragraphs[0].space_after = Pt(14.4)  # Line spacing 1.2 * 12pt

    bullet_points = [
        "The Transformer model generalizes well to English constituency parsing.",
        "Competitive performance observed in both supervised and semi-supervised settings.",
        "Demonstrates the model's versatility beyond translation tasks."
    ]

    for point in bullet_points:
        p = bullet_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(163, 185, 204)
        p.font.name = 'Arial'
        p.space_after = Pt(14.4)  # Line spacing 1.2 * 12pt

    # Add table
    rows, cols = 2, 3  # Example table size
    table_box = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.5), Inches(9), Inches(3))
    table = table_box.table

    # Set table content (example content)
    table.cell(0, 0).text = "Model"
    table.cell(0, 1).text = "Accuracy"
    table.cell(0, 2).text = "F1 Score"
    table.cell(1, 0).text = "Transformer"
    table.cell(1, 1).text = "90%"
    table.cell(1, 2).text = "0.88"

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Discuss the generalization capabilities of the Transformer model to other tasks."

    # Save the presentation
    prs.save(output_pptx_path)

if __name__ == "__main__":
    output_pptx_path = "output/0116_2040/result/slide_08/slide.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    create_presentation(output_pptx_path)