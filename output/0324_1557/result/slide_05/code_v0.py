import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def add_slide(prs):
    slide_layout = prs.slide_layouts[5]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Header Bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Footer Bar
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.5), Inches(10), Inches(1.125)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)

    footer_text_box = footer_bar.text_frame
    footer_text_box.text = "Slide Number / Presenter Name"
    footer_text_box.paragraphs[0].font.name = "Microsoft YaHei"
    footer_text_box.paragraphs[0].font.size = Pt(14)
    footer_text_box.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    footer_text_box.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Horizontal Separator Line
    separator_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.125), Inches(9), Inches(0.01)
    )
    separator_line.fill.solid()
    separator_line.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.625)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Technical Details: Self-Attention and Positional Encoding"
    title_frame.paragraphs[0].font.name = "Microsoft YaHei"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Body Content
    body_content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5)
    )
    body_content_frame = body_content_box.text_frame
    body_content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_content_frame.word_wrap = True
    body_content_frame.fill.solid()
    body_content_frame.fill.fore_color.rgb = RGBColor(230, 240, 250)

    bullet_points = [
        "Self-attention enables the model to capture dependencies between tokens, regardless of their distance in the sequence.",
        "Positional encoding uses sine and cosine functions to represent the position of tokens.",
        "These components work together to model sequence order and relationships effectively."
    ]

    for point in bullet_points:
        p = body_content_frame.add_paragraph()
        p.text = point
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.LEFT

    # Figure
    image_path = "output/0324_1557/images/_page_12_Figure_1.jpeg"
    figure = slide.shapes.add_picture(
        image_path, Inches(6.5), Inches(1.5), width=Inches(3)
    )

    # Figure Caption
    caption_box = slide.shapes.add_textbox(
        Inches(6.5), Inches(6), Inches(3), Inches(0.5)
    )
    caption_frame = caption_box.text_frame
    caption_frame.text = "An example of the attention mechanism following long-distance dependencies in the encoder self-attention in layer 5 of 6."
    caption_frame.paragraphs[0].font.name = "Microsoft YaHei"
    caption_frame.paragraphs[0].font.size = Pt(14)
    caption_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker Notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Dive deeper into self-attention and positional encoding, using a figure to illustrate long-distance dependencies."

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    output_pptx_path = "output/0324_1557/result/slide_05/slide.pptx"
    prs.save(output_pptx_path)