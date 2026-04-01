import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    # Add a new slide with Title and Content layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Set the background color of the slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray

    # Add header bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)  # Deep blue
    header_bar.line.fill.background()  # Remove border

    # Add footer bar
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.625 - 1.125), Inches(10), Inches(1.125)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
    footer_bar.line.fill.background()  # Remove border

    # Add footer text
    footer_text_box = footer_bar.text_frame
    footer_text_box.text = "Presenter Name | Slide Number"
    footer_text_box.word_wrap = True
    footer_text_box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    footer_text_box.paragraphs[0].font.name = "Microsoft YaHei"
    footer_text_box.paragraphs[0].font.size = Pt(14)
    footer_text_box.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)  # Dark text
    footer_text_box.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Add horizontal separator line below header bar
    separator_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.125), Inches(9), Inches(0.01)
    )
    separator_line.fill.solid()
    separator_line.fill.fore_color.rgb = RGBColor(0, 63, 135)  # Deep blue
    separator_line.line.fill.background()  # Remove border

    # Add title text in the header bar
    title_text_box = header_bar.text_frame
    title_text_box.text = "Technical Details: Self-Attention and Positional Encoding"
    title_text_box.word_wrap = True
    title_text_box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_paragraph = title_text_box.paragraphs[0]
    title_paragraph.font.name = "Microsoft YaHei"
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Add body content area
    body_content_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5)
    )
    body_content_box.fill.solid()
    body_content_box.fill.fore_color.rgb = RGBColor(230, 240, 250)  # Light blue
    body_content_box.line.fill.background()  # Remove border

    # Add bullet points to the body content
    body_text_frame = body_content_box.text_frame
    body_text_frame.word_wrap = True
    body_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_points = [
        "Self-attention enables the model to capture dependencies between tokens, regardless of their distance in the sequence.",
        "Positional encoding uses sine and cosine functions to represent the position of tokens.",
        "These components work together to model sequence order and relationships effectively."
    ]
    for point in bullet_points:
        p = body_text_frame.add_paragraph()
        p.text = point
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)  # Dark text
        p.level = 0

    # Add figure to the right content area
    image_path = "output/0324_1557/images/_page_12_Figure_1.jpeg"
    figure = slide.shapes.add_picture(
        image_path, Inches(6.5), Inches(1.5), width=Inches(3)
    )

    # Add caption below the figure
    caption_box = slide.shapes.add_textbox(
        Inches(6.5), Inches(6), Inches(3), Inches(0.5)
    )
    caption_text_frame = caption_box.text_frame
    caption_text_frame.text = "An example of the attention mechanism following long-distance dependencies in the encoder self-attention in layer 5 of 6."
    caption_text_frame.word_wrap = True
    caption_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    caption_paragraph = caption_text_frame.paragraphs[0]
    caption_paragraph.font.name = "Microsoft YaHei"
    caption_paragraph.font.size = Pt(14)
    caption_paragraph.font.color.rgb = RGBColor(51, 51, 51)  # Dark text
    caption_paragraph.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Dive deeper into self-attention and positional encoding, using a figure to illustrate long-distance dependencies."

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_05/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)