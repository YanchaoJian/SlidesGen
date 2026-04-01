import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def add_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
    background.line.width = Pt(0)  # No border

    # Add top title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)  # Deep blue
    title_bar.line.width = Pt(0)  # No border

    # Add footer area
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.625 - 1.125), Inches(10), Inches(1.125)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
    footer.line.width = Pt(0)  # No border

    # Add footer text
    footer_text_box = slide.shapes.add_textbox(Inches(8), Inches(5.625 - 0.75), Inches(2), Inches(0.375))
    footer_text_frame = footer_text_box.text_frame
    footer_text_frame.word_wrap = True
    footer_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    footer_text = footer_text_frame.add_paragraph()
    footer_text.text = "Slide 10 | Academic Blue Professional Theme"
    footer_text.font.name = "Microsoft YaHei"
    footer_text.font.size = Pt(14)
    footer_text.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    footer_text.alignment = PP_ALIGN.RIGHT

    # Add title text
    title_text_box = slide.shapes.add_textbox(Inches(0), Inches(0.375), Inches(10), Inches(0.375))
    title_text_frame = title_text_box.text_frame
    title_text_frame.word_wrap = True
    title_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_text = title_text_frame.add_paragraph()
    title_text.text = "Questions & Discussion"
    title_text.font.name = "Microsoft YaHei"
    title_text.font.size = Pt(32)
    title_text.font.bold = True
    title_text.font.color.rgb = RGBColor(255, 255, 255)  # White
    title_text.alignment = PP_ALIGN.CENTER

    # Add body content
    body_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.875))
    body_text_box.fill.solid()
    body_text_box.fill.fore_color.rgb = RGBColor(230, 240, 250)  # Light blue
    body_text_box.line.width = Pt(0)  # No border
    body_text_frame = body_text_box.text_frame
    body_text_frame.word_wrap = True
    body_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_text_1 = body_text_frame.add_paragraph()
    body_text_1.text = "Thank you for your attention!"
    body_text_1.font.name = "Microsoft YaHei"
    body_text_1.font.size = Pt(18)
    body_text_1.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    body_text_1.alignment = PP_ALIGN.LEFT
    body_text_2 = body_text_frame.add_paragraph()
    body_text_2.text = "Questions and feedback are welcome."
    body_text_2.font.name = "Microsoft YaHei"
    body_text_2.font.size = Pt(18)
    body_text_2.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    body_text_2.alignment = PP_ALIGN.LEFT

    # Add presenter notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Encourage audience engagement and discussion."

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_10/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)