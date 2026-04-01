import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set slide background color to light gray (245, 245, 245)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 245)
    background.line.fill.background()  # Remove border

    # Header Bar
    header_height = Inches(1.125)  # 20% of slide height
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, header_height)
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)  # Deep Blue
    header_bar.line.fill.background()  # Remove border

    # Thin horizontal line below header bar
    line_top = header_height
    line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(0.5), line_top, Inches(9.5), line_top)
    line.line.color.rgb = RGBColor(0, 63, 135)  # Deep Blue
    line.line.width = Pt(1)  # 1pt thickness

    # Footer Bar
    footer_top = prs.slide_height - header_height
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_top, prs.slide_width, header_height)
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light Gray
    footer_bar.line.fill.background()  # Remove border

    # Footer Text
    footer_text_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.375), Inches(1.0), Inches(0.5))
    footer_text_frame = footer_text_box.text_frame
    footer_text_frame.text = "Academic Blue Professional Theme"
    footer_text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    footer_text_frame.paragraphs[0].font.size = Pt(14)
    footer_text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)  # Dark Text
    footer_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Title in Header Bar
    title_text_box = slide.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, header_height)
    title_text_frame = title_text_box.text_frame
    title_text_frame.text = "Solution Overview: The Transformer Model"
    title_text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    title_text_frame.paragraphs[0].font.size = Pt(32)
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Body Text Box
    body_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(3.375))
    body_text_frame = body_text_box.text_frame
    body_text_frame.word_wrap = True
    body_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_text_frame.text = (
        "The Transformer introduces a novel architecture based entirely on attention mechanisms, "
        "eliminating recurrence and convolution.\n\n"
        "This approach enhances parallelization, reduces training time, and improves translation quality.\n\n"
        "Key innovations include multi-head self-attention and positional encoding to capture sequence order."
    )
    for paragraph in body_text_frame.paragraphs:
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)  # Dark Text
        paragraph.alignment = PP_ALIGN.LEFT

    # Image on the right side
    image_path = "output/0324_1557/images/_page_2_Figure_0.jpeg"
    image_left = Inches(5.5)
    image_top = Inches(1.5)
    image_width = Inches(3.5)
    image_height = Inches(3.375)
    slide.shapes.add_picture(image_path, image_left, image_top, width=image_width, height=image_height)

    # Caption below the image
    caption_text_box = slide.shapes.add_textbox(Inches(5.5), Inches(5.0), Inches(3.5), Inches(0.5))
    caption_text_frame = caption_text_box.text_frame
    caption_text_frame.text = "The Transformer - model architecture."
    caption_text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    caption_text_frame.paragraphs[0].font.size = Pt(14)
    caption_text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)  # Dark Text
    caption_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Presenter Notes
    slide.notes_slide.notes_text_frame.text = (
        "Provide a high-level overview of the Transformer model and its key innovations, supported by the architecture diagram."
    )

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_03/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)