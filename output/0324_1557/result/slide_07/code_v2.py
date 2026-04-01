import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set slide background color to light gray
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Static Decoration Layer
    # Top title bar
    title_bar_height = Inches(1.125)  # 20% of slide height
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=0, 
        top=0, 
        width=prs.slide_width, 
        height=title_bar_height
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Bottom footer area
    footer_height = Inches(1.125)  # 20% of slide height
    footer_top = prs.slide_height - footer_height
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=0, 
        top=footer_top, 
        width=prs.slide_width, 
        height=footer_height
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Thin horizontal line below the title bar
    line_top = title_bar_height
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, 
        start_x=0, 
        start_y=line_top, 
        end_x=prs.slide_width, 
        end_y=line_top
    )
    line.line.color.rgb = RGBColor(0, 63, 135)
    line.line.width = Pt(1)

    # Dynamic Layout Layer
    # Title
    title_top = Inches(0.5)
    title = slide.shapes.add_textbox(
        left=Inches(0.5), 
        top=title_top, 
        width=prs.slide_width - Inches(1), 
        height=Inches(1)
    )
    text_frame = title.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = text_frame.add_paragraph()
    p.text = "Results: Translation Quality Comparison"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 63, 135)

    # Body/List
    body_top = Inches(1.75)
    body = slide.shapes.add_textbox(
        left=Inches(0.5), 
        top=body_top, 
        width=Inches(8), 
        height=Inches(2)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(230, 240, 250)
    text_frame = body.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_points = [
        "The Transformer achieves 28.4 BLEU on English-to-German, improving over the existing best results by over 2 BLEU.",
        "On English-to-French, the Transformer establishes a new single-model state-of-the-art BLEU score of 41.8.",
        "Training time was significantly reduced compared to previous models."
    ]
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.level = 0

    # Table
    table_top = Inches(3.75)
    table = slide.shapes.add_table(
        rows=2, 
        cols=2, 
        left=Inches(0.5), 
        top=table_top, 
        width=Inches(8), 
        height=Inches(1.5)
    ).table
    table.first_row = True
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    table.rows[0].height = Inches(0.75)
    table.rows[1].height = Inches(0.75)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "BLEU Score"
    table.cell(1, 1).text = "28.4 (English-to-German), 41.8 (English-to-French)"
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 240, 250)

    # Speaker Notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Highlight the superior BLEU scores and reduced training costs achieved by the Transformer."

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_07/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)