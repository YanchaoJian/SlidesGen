import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set slide background color to light gray
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Static Decoration Layer
    # Top title bar
    title_bar_height = Inches(1.125)  # 20% of slide height
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=0, 
        top=0, 
        width=prs.slide_width, 
        height=title_bar_height
    ).fill.solid().fore_color.rgb = RGBColor(0, 63, 135)

    # Bottom footer area
    footer_top = prs.slide_height - title_bar_height
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=0, 
        top=footer_top, 
        width=prs.slide_width, 
        height=title_bar_height
    ).fill.solid().fore_color.rgb = RGBColor(245, 245, 245)

    # Thin horizontal line below the title bar
    line_top = title_bar_height
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, 
        start_x=0, 
        start_y=line_top, 
        end_x=prs.slide_width, 
        end_y=line_top
    )
    line.line.color.rgb = RGBColor(0, 63, 135)
    line.line.width = Pt(1)

    # Dynamic Layout Layer
    # Title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = prs.slide_width - Inches(1)
    title_height = Inches(1)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.text = "Results: Translation Quality Comparison"
    title_frame.paragraphs[0].font.name = "Microsoft YaHei"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 63, 135)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Body/List
    body_left = Inches(0.5)
    body_top = Inches(1.75)
    body_width = Inches(8)
    body_height = Inches(2)
    body_box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_box.fill.solid()
    body_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
    body_frame.text = (
        "• The Transformer achieves 28.4 BLEU on English-to-German, improving over the existing best results by over 2 BLEU.\n"
        "• On English-to-French, the Transformer establishes a new single-model state-of-the-art BLEU score of 41.8.\n"
        "• Training time was significantly reduced compared to previous models."
    )
    for paragraph in body_frame.paragraphs:
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)

    # Table
    table_top = Inches(3.75)
    table_left = Inches(0.5)
    table_width = Inches(8)
    table_height = Inches(1.5)
    table_box = slide.shapes.add_textbox(table_left, table_top, table_width, table_height)
    table_frame = table_box.text_frame
    table_frame.word_wrap = True
    table_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    table_box.fill.solid()
    table_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
    table_frame.text = (
        "The Transformer achieves better BLEU scores than previous state-of-the-art models on the English-to-German and "
        "English-to-French newstest2014 tests at a fraction of the training cost."
    )
    table_frame.paragraphs[0].font.name = "Microsoft YaHei"
    table_frame.paragraphs[0].font.size = Pt(18)
    table_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

    # Speaker Notes
    slide.notes_slide.notes_text_frame.text = (
        "Highlight the superior BLEU scores and reduced training costs achieved by the Transformer."
    )

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_07/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)