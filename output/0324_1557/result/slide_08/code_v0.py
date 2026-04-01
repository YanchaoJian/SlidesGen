import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def add_slide(prs):
    # Add a slide with "Title and Content" layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Set slide background color to light gray (RGB: 245, 245, 245)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Add top title bar (deep blue rectangle)
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)
    title_bar.line.fill.background()  # Remove border
    
    # Add footer bar (light gray rectangle)
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.625 - 1.125), Inches(10), Inches(1.125)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)
    footer_bar.line.fill.background()  # Remove border
    
    # Add footer text
    footer_text_box = slide.shapes.add_textbox(Inches(8), Inches(5.625 - 0.875), Inches(2), Inches(0.5))
    footer_text_frame = footer_text_box.text_frame
    footer_text_frame.word_wrap = True
    footer_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    footer_text = footer_text_frame.add_paragraph()
    footer_text.text = "Slide 8 | Academic Blue Professional Theme"
    footer_text.font.name = "Microsoft YaHei"
    footer_text.font.size = Pt(14)
    footer_text.font.color.rgb = RGBColor(51, 51, 51)
    footer_text.alignment = PP_ALIGN.RIGHT
    
    # Add title text in the top title bar
    title_text_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1.125))
    title_text_frame = title_text_box.text_frame
    title_text_frame.word_wrap = True
    title_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_text = title_text_frame.add_paragraph()
    title_text.text = "Analysis: Attention Mechanism Insights"
    title_text.font.name = "Microsoft YaHei"
    title_text.font.size = Pt(32)
    title_text.font.bold = True
    title_text.font.color.rgb = RGBColor(255, 255, 255)
    title_text.alignment = PP_ALIGN.CENTER
    
    # Add body text box (light blue rectangle)
    body_text_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5), Inches(4.875)
    )
    body_text_box.fill.solid()
    body_text_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
    body_text_box.line.fill.background()  # Remove border
    body_text_frame = body_text_box.text_frame
    body_text_frame.word_wrap = True
    body_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_text = body_text_frame.add_paragraph()
    body_text.text = (
        "Attention heads capture long-distance dependencies, crucial for understanding complex sentence structures.\n"
        "Some attention heads specialize in tasks like anaphora resolution and syntactic structure modeling.\n"
        "These insights demonstrate the versatility and effectiveness of the attention mechanism."
    )
    body_text.font.name = "Microsoft YaHei"
    body_text.font.size = Pt(18)
    body_text.font.color.rgb = RGBColor(51, 51, 51)
    body_text.alignment = PP_ALIGN.LEFT
    
    # Add image on the right side
    image_path = "output/0324_1557/images/_page_13_Figure_0.jpeg"
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(5.75), Inches(1.5), width=Inches(3.75), height=Inches(4.875))
    
    # Add image caption below the image
    caption_text_box = slide.shapes.add_textbox(Inches(5.75), Inches(6.375), Inches(3.75), Inches(0.5))
    caption_text_frame = caption_text_box.text_frame
    caption_text_frame.word_wrap = True
    caption_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    caption_text = caption_text_frame.add_paragraph()
    caption_text.text = "Two attention heads, also in layer 5 of 6, apparently involved in anaphora resolution."
    caption_text.font.name = "Microsoft YaHei"
    caption_text.font.size = Pt(14)
    caption_text.font.color.rgb = RGBColor(51, 51, 51)
    caption_text.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = (
        "Discuss the qualitative insights gained from analyzing attention heads, supported by visual examples."
    )

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_08/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)