import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # [1. Global Settings] Set background color to light gray
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # [2. Static Decoration Layer]
    # Top title area
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)
    title_bar.line.color.rgb = RGBColor(0, 63, 135)  # Ensure no border

    # Thin horizontal line below the title area
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Inches(0), Inches(1.125), Inches(10), Inches(1.125)
    )
    line.line.color.rgb = RGBColor(0, 63, 135)
    line.line.width = Pt(1)

    # Bottom area
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.5), Inches(10), Inches(1.125)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)
    bottom_bar.line.color.rgb = RGBColor(245, 245, 245)  # Ensure no border

    # [3. Dynamic Layout Layer]
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.text = "Methodology: Attention Mechanisms"
    title_frame.paragraphs[0].font.name = "Microsoft YaHei"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 63, 135)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Body/List
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(3.375))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_frame.paragraphs[0].font.name = "Microsoft YaHei"
    body_frame.paragraphs[0].font.size = Pt(18)
    body_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    body_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add bullet points
    bullet_points = [
        "The Transformer uses scaled dot-product attention to compute attention weights efficiently.",
        "Multi-head attention allows the model to focus on different parts of the input sequence simultaneously.",
        "Positional encoding helps the model understand the order of tokens in a sequence."
    ]
    for point in bullet_points:
        p = body_frame.add_paragraph()
        p.text = point
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.level = 0

    # Visual Assets
    image_path = "output/0324_1557/images/_page_3_Figure_0.jpeg"
    if os.path.exists(image_path):
        image = slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), height=Inches(3.375))
    else:
        print(f"Image not found at path: {image_path}")

    # Caption below the image
    caption_box = slide.shapes.add_textbox(Inches(5.5), Inches(4), Inches(4.5), Inches(0.5))
    caption_frame = caption_box.text_frame
    caption_frame.word_wrap = True
    caption_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    caption_frame.text = "(left) Scaled Dot-Product Attention. (right) Multi-Head Attention consists of several attention layers running in parallel."
    caption_frame.paragraphs[0].font.name = "Microsoft YaHei"
    caption_frame.paragraphs[0].font.size = Pt(14)
    caption_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    caption_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # [4. Notes]
    slide.notes_slide.notes_text_frame.text = "Explain the core attention mechanisms used in the Transformer, supported by visual illustrations."

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_04/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)