import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set background color to light gray
    slide_background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    slide_background.fill.solid()
    slide_background.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Top title area with deep blue background
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.125)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Thin horizontal line below the title area
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.125), prs.slide_width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Bottom area with light gray background
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(1.125), prs.slide_width, Inches(1.125)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    title_frame = title_box.text_frame
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.word_wrap = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Methodology: Attention Mechanisms"
    title_run.font.bold = True
    title_run.font.size = Pt(32)
    title_run.font.name = 'Microsoft YaHei'
    title_run.font.color.rgb = RGBColor(0, 63, 135)

    # Bullet points text box
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(3.375))
    text_frame = text_box.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    bullet_points = [
        "The Transformer uses scaled dot-product attention to compute attention weights efficiently.",
        "Multi-head attention allows the model to focus on different parts of the input sequence simultaneously.",
        "Positional encoding helps the model understand the order of tokens in a sequence."
    ]

    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(10)

    # Image on the right side
    image_path = "output/0324_1557/images/_page_3_Figure_0.jpeg"
    slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), height=Inches(3.375))

    # Image caption
    caption_box = slide.shapes.add_textbox(Inches(5.5), Inches(4.875), Inches(4), Inches(0.5))
    caption_frame = caption_box.text_frame
    caption_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    caption_frame.word_wrap = True
    caption_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    caption_run = caption_frame.paragraphs[0].add_run()
    caption_run.text = "(left) Scaled Dot-Product Attention. (right) Multi-Head Attention consists of several attention layers running in parallel."
    caption_run.font.size = Pt(12)
    caption_run.font.name = 'Microsoft YaHei'
    caption_run.font.color.rgb = RGBColor(51, 51, 51)

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Explain the core attention mechanisms used in the Transformer, supported by visual illustrations."

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    output_pptx_path = "output/0324_1557/result/slide_04/slide.pptx"
    prs.save(output_pptx_path)