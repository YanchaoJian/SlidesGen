import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

def add_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set slide background color to light gray
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 245)
    background.line.color.rgb = RGBColor(245, 245, 245)  # Remove border

    # Add header line at the top
    header_line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Inches(0), Inches(0.5), prs.slide_width, Inches(0.5)
    )
    header_line.line.color.rgb = RGBColor(0, 63, 135)
    header_line.line.width = Pt(1)

    # Add footer area at the bottom
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(1.125), prs.slide_width, Inches(1.125)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(245, 245, 245)
    footer.line.color.rgb = RGBColor(245, 245, 245)  # Remove border

    # Add title text box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Background: Sequence Transduction Models in AI"
    title_run.font.name = "Microsoft YaHei"
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 63, 135)

    # Add bullet points text box
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.625), Inches(8), Inches(3.375))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
    content_box.line.color.rgb = RGBColor(230, 240, 250)  # Remove border

    # Add bullet points
    bullet_points = [
        "Sequence transduction models are foundational for tasks like language modeling, machine translation, and parsing.",
        "Traditional models rely on recurrent neural networks (RNNs) and convolutional neural networks (CNNs), which are computationally intensive.",
        "Attention mechanisms have emerged as a powerful tool to model dependencies in sequences, enabling more efficient computation."
    ]
    for point in bullet_points:
        p = content_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = (
        "Introduce the audience to the importance of sequence transduction models "
        "and the role of attention mechanisms in advancing the field."
    )

if __name__ == "__main__":
    output_pptx_path = "output/0324_1557/result/slide_01/slide.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)