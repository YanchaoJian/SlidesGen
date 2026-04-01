import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set background color to light gray
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Draw top title area with deep blue background
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Add title text
    title_text_box = slide.shapes.add_textbox(
        Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    title_text_frame = title_text_box.text_frame
    title_text_frame.text = "Experimental Setup: Datasets and Metrics"
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].font.size = Pt(32)
    title_text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Draw a thin horizontal line below the title area
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.125), Inches(10), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Draw bottom area with light gray background
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.125), Inches(10), Inches(0.5)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Add footer text
    footer_text_box = slide.shapes.add_textbox(
        Inches(0), Inches(5.125), Inches(10), Inches(0.5)
    )
    footer_text_frame = footer_text_box.text_frame
    footer_text_frame.text = "Footer text"
    footer_text_frame.paragraphs[0].font.size = Pt(12)
    footer_text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    footer_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Add body text box for bullet points
    body_text_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.625), Inches(9.0), Inches(3.375)
    )
    body_text_frame = body_text_box.text_frame
    body_text_frame.word_wrap = True
    body_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Add bullet points
    bullet_points = [
        "Experiments were conducted on WMT 2014 English-to-German and English-to-French translation tasks.",
        "BLEU scores were used as the primary evaluation metric to measure translation quality.",
        "Baseline methods included state-of-the-art RNN and CNN models."
    ]
    for point in bullet_points:
        p = body_text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(10)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Provide an overview of the experimental setup, including datasets, metrics, and baseline methods."

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    output_pptx_path = 'output/0324_1557/result/slide_06/slide.pptx'
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)