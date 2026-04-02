import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set slide background color
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Header Bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.125)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 63, 135)

    # Title in Header Bar
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.375), Inches(9), Inches(0.375)
    )
    title_frame = title_box.text_frame
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.word_wrap = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Problem: Limitations of Traditional Models"
    title_run.font.bold = True
    title_run.font.size = Pt(32)
    title_run.font.color.rgb = RGBColor(255, 255, 255)

    # Footer Bar
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.375), Inches(10), Inches(1.125)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Footer Text
    footer_text_box = slide.shapes.add_textbox(
        Inches(8), Inches(6.625), Inches(2), Inches(0.5)
    )
    footer_text_frame = footer_text_box.text_frame
    footer_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    footer_text_frame.word_wrap = True
    footer_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    footer_run = footer_text_frame.paragraphs[0].add_run()
    footer_run.text = "Academic Blue Professional Theme"
    footer_run.font.size = Pt(14)
    footer_run.font.color.rgb = RGBColor(51, 51, 51)

    # Body Content
    body_content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(4.875)
    )
    body_content_box.fill.solid()
    body_content_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
    body_content_frame = body_content_box.text_frame
    body_content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body_content_frame.word_wrap = True
    body_content_frame.margin_left = Inches(0.5)
    body_content_frame.margin_top = Inches(0.5)
    body_content_frame.margin_right = Inches(0.5)
    body_content_frame.margin_bottom = Inches(0.5)

    body_text = (
        "RNNs and CNNs face challenges in parallelization due to their sequential nature, "
        "which affects training efficiency and scalability.\n"
        "These models struggle with long-range dependencies, requiring significant computational resources.\n"
        "A more efficient and scalable approach is needed to address these limitations."
    )
    body_paragraph = body_content_frame.add_paragraph()
    body_paragraph.text = body_text
    body_paragraph.font.size = Pt(18)
    body_paragraph.font.color.rgb = RGBColor(51, 51, 51)

    # Presenter Notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = (
        "Highlight the specific challenges faced by traditional models and set the stage for the proposed solution."
    )

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_slide(prs)
    output_pptx_path = "output/0324_1557/result/slide_02/slide.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)