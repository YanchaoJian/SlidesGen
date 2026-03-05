from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants for style
BACKGROUND_COLOR = RGBColor(255, 255, 255)
HEADER_BAR_COLOR = RGBColor(0, 112, 192)
TITLE_COLOR = RGBColor(0, 0, 0)
BODY_TEXT_COLOR = RGBColor(169, 169, 169)

# Error handling for file path
output_path = 'output/1230_1234/result/slide_06/slide.pptx'
output_dir = os.path.dirname(output_path)
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Create presentation and slide
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a predefined layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BACKGROUND_COLOR

# Add master shapes
# Add header bar rectangle
header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(10.0), Inches(0.5)
)
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = HEADER_BAR_COLOR

# Add vertical line using a shape instead of add_line
line = slide.shapes.add_shape(
    MSO_SHAPE.LINE, Inches(1.0), Inches(0.0), Inches(1.0), Inches(0.5)
)
line.line.color.rgb = HEADER_BAR_COLOR

# Add title box
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(8.0), Inches(1.0))
title = title_box.text_frame
title.text = 'Experimental Setup'
title.paragraphs[0].font.size = Pt(28)
title.paragraphs[0].font.name = 'Arial'
title.paragraphs[0].font.color.rgb = TITLE_COLOR
title.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add body text box
body_text_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.125))
body_text_frame = body_text_box.text_frame

# Add bullet points
bullet_points = [
    'Experiments conducted on WMT 2014 English-to-German and English-to-French translation tasks.',
    'Evaluation metrics include BLEU scores.',
    'Baseline methods include state-of-the-art models like GNMT and ConvS2S.'
]

for point in bullet_points:
    paragraph = body_text_frame.add_paragraph()
    paragraph.text = point
    paragraph.font.size = Pt(18)
    paragraph.font.name = 'Arial'
    paragraph.font.color.rgb = BODY_TEXT_COLOR

# Save presentation to specified path
presentation.save(output_path)