from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a presentation object
prs = Presentation()

# Add a slide with a 'Title and Content' layout
slide_layout = prs.slide_layouts[5]  # Using 'Title and Content' layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw the top header bar
header_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.0), Inches(0.0), Inches(10.0), Inches(0.5)
)
header_fill = header_shape.fill
header_fill.solid()
header_fill.fore_color.rgb = RGBColor(0, 112, 192)

# Draw the vertical line
line = slide.shapes.add_line(
    Inches(1.0), Inches(0.0), Inches(1.0), Inches(0.5)
)
line.line.color.rgb = RGBColor(0, 112, 192)
line.line.width = Pt(2.0)

# Create title box
title_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(0.5), Inches(8.0), Inches(1.0)
)
title_frame = title_box.text_frame
title_frame.text = "Experimental Setup"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Create a text box for bullet points
content_box = slide.shapes.add_textbox(
    Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.125)
)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add bullet points
bullet_points = [
    "• Experiments conducted on WMT 2014 English-to-German and English-to-French translation tasks.",
    "• Evaluation metrics include BLEU scores.",
    "• Baseline methods include state-of-the-art models like GNMT and ConvS2S."
]

for point in bullet_points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)

# Ensure the first paragraph is not a bullet point
content_frame.paragraphs[0].level = 0

# Save the presentation
output_path = 'output/1230_1234/result/slide_06/slide.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)