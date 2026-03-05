from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create a presentation object
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[5]  # 5 is typically a blank layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add a rectangle shape (top_header_bar)
left = Inches(0.0)
top = Inches(0.0)
width = Inches(10.0)
height = Inches(0.5)
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Add a vertical line (vertical_line)
left = Inches(1.0)
top = Inches(0.0)
width = Inches(0.0)
height = Inches(0.5)
line = slide.shapes.add_line(left, top, left, top + height)
line.line.color.rgb = RGBColor(0, 112, 192)

# Add title box
left = Inches(1.5)
top = Inches(0.5)
width = Inches(8.0)
height = Inches(1.0)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = 'Conclusion and Future Directions'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add body text box
left = Inches(1.0)
top = Inches(1.5)
width = Inches(8.0)
height = Inches(4.125)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame

# Add bullet points
bullet_points = [
    'The Transformer model offers a more efficient and scalable approach to sequence transduction.',
    'Achieved state-of-the-art results in machine translation tasks.',
    'Future work includes exploring other applications and further optimizing the model.'
]

for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)

# Ensure the first paragraph is not a bullet point
text_frame.paragraphs[0].level = 0

# Save the presentation
output_path = 'output/1230_1234/result/slide_09/slide.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)