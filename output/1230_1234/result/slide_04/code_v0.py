import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt

# Create a presentation object
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw the top header bar
header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(10.0), Inches(0.5)
)
header_fill = header_bar.fill
header_fill.solid()
header_fill.fore_color.rgb = RGBColor(0, 112, 192)

# Draw the vertical line
line = slide.shapes.add_shape(
    MSO_SHAPE.LINE, Inches(1.0), Inches(0.0), Inches(1.0), Inches(0.5)
)
line.line.color.rgb = RGBColor(0, 112, 192)

# Add title
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(8.0), Inches(1.0))
title_frame = title_box.text_frame
title_frame.text = 'Methodology: Attention Mechanisms'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add body text with bullet points
body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.8), Inches(3.0))
body_frame = body_box.text_frame
body_frame.word_wrap = True

bullet_points = [
    'The Transformer uses a multi-head attention mechanism.',
    'Includes scaled dot-product attention and positional encoding.',
    'Handles sequence order without recurrence.'
]

for point in bullet_points:
    p = body_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.space_after = Pt(10)

# Add image
image_path = 'output/1230_1234/images/_page_3_Figure_1.jpeg'
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(6.0), Inches(1.6), Inches(3.0), Inches(2.5))
else:
    print(f"Warning: Image file '{image_path}' not found.")

# Render LaTeX equation using matplotlib
latex_string = r'Attention(Q, K, V) = \operatorname{softmax}\left(\frac{QK^T}{\sqrt{d_k}}\right)V'
plt.figure(figsize=(4, 0.5))
plt.text(0.5, 0.5, f"${latex_string}$", fontsize=20, ha='center', va='center')
plt.axis('off')
equation_image_path = 'temp_equation.png'
plt.savefig(equation_image_path, bbox_inches='tight', pad_inches=0, dpi=300)
plt.close()

# Add equation image
if os.path.exists(equation_image_path):
    slide.shapes.add_picture(equation_image_path, Inches(1.0), Inches(4.2))
else:
    print(f"Warning: Equation image file '{equation_image_path}' not found.")

# Save the presentation
output_path = 'output/1230_1234/result/slide_04/slide.pptx'
prs.save(output_path)
print(f"Presentation saved to '{output_path}'")