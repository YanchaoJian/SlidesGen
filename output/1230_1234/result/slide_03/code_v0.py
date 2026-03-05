import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Use blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw a rectangle (top header bar)
header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.0), Inches(0.0), Inches(10.0), Inches(0.5)
)
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Draw a vertical line
line = slide.shapes.add_shape(
    MSO_SHAPE.LINE,
    Inches(1.0), Inches(0.0), Inches(1.0), Inches(0.5)
)
line.line.color.rgb = RGBColor(0, 112, 192)

# Create title box
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(8.0), Inches(1.0))
title_frame = title_box.text_frame
title_frame.text = 'Our Core Contribution: The Transformer Model'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Create body text box
body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(8.0), Inches(2.0))
body_frame = body_box.text_frame
body_frame.word_wrap = True

# Add bullet points
bullet_points = [
    'Introduced a novel architecture based solely on attention mechanisms.',
    'Eliminates the need for recurrence and convolution.',
    'Allows for greater parallelization and improved translation quality.'
]

for point in bullet_points:
    p = body_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.space_after = Pt(10)

# Insert image
image_path = 'output/1230_1234/images/_page_2_Figure_0.jpeg'
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(1.0), Inches(3.8), width=Inches(8.0), height=Inches(1.825))
else:
    print(f"Warning: Image file '{image_path}' not found.")

# Save the presentation
output_path = 'output/1230_1234/result/slide_03/slide.pptx'
presentation.save(output_path)