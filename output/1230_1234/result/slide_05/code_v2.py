import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt

# Define the paths
slide_save_path = 'output/1230_1234/result/slide_05/slide.pptx'
image_path = 'output/1230_1234/images/_page_3_Figure_2.jpeg'

# Create a new presentation object
presentation = Presentation()

# Add a new slide layout
slide_layout = presentation.slide_layouts[5]  # Using a blank layout here
slide = presentation.slides.add_slide(slide_layout)

# Set background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw the "Top Header Bar"
top_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.0), Inches(0.0),
    Inches(10.0), Inches(0.5)
)
bar_fill = top_bar.fill
bar_fill.solid()
bar_fill.fore_color.rgb = RGBColor(0, 112, 192)

# Draw the "Vertical Line"
vertical_line = slide.shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,  # Corrected shape type for a line
    Inches(1.0), Inches(0.0),
    Inches(1.0), Inches(0.5)  # Width is zero for a vertical line
)
line_fill = vertical_line.line
line_fill.color.rgb = RGBColor(0, 112, 192)

# Add Title
title_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(0.5),
    Inches(8.0), Inches(1.0)
)
title_frame = title_box.text_frame
title_para = title_frame.add_paragraph()
title_para.text = "Technical Approach: Multi-Head Attention"
title_para.font.size = Inches(0.39)  # 28pt is around 39 points and 1pt is ~1/72 inches
title_para.font.name = 'Arial'
title_para.font.color.rgb = RGBColor(0, 0, 0)
title_para.alignment = PP_ALIGN.LEFT

# Add Body Text
body_box = slide.shapes.add_textbox(
    Inches(1.0), Inches(1.5),
    Inches(8.0), Inches(1.5)
)
body_frame = body_box.text_frame
body_texts = [
    "Multi-head attention allows the model to attend to different parts of the input sequence simultaneously.",
    "Each head performs scaled dot-product attention.",
    "Results are concatenated and linearly transformed."
]
for text in body_texts:
    body_para = body_frame.add_paragraph()
    body_para.text = text
    body_para.font.size = Inches(0.25)  # 18pt is around 18/72 inches
    body_para.font.name = 'Arial'
    body_para.font.color.rgb = RGBColor(169, 169, 169)
body_frame.text_anchor = PP_ALIGN.LEFT

# Add Image
if os.path.exists(image_path):
    slide.shapes.add_picture(
        image_path,
        Inches(1.0), Inches(3.0),
        width=Inches(4.0),
        height=Inches(2.625)
    )
else:
    print(f"Warning: Image file '{image_path}' not found.")

# Render LaTeX Equation using matplotlib
latex_string = r"""
\mathrm{MultiHead}(Q, K, V) = \mathrm{Concat}(\mathrm{head}_1, \dots, \mathrm{head}_h) W^O \\
\mathrm{where\ } \mathrm{head}_i = \mathrm{Attention}(QW_i^Q, KW_i^K, VW_i^V)
"""
temp_path = 'temp_equation.png'
plt.figure(figsize=(6, 2))
plt.text(0.1, 0.5, f"${latex_string}$", fontsize=12, va='center')
plt.axis('off')
plt.savefig(temp_path, bbox_inches='tight', pad_inches=0, dpi=100)

# Add Equation Image to slide
if os.path.exists(temp_path):
    slide.shapes.add_picture(
        temp_path,
        Inches(5.5), Inches(3.0)
    )
else:
    print(f"Warning: Equation image file '{temp_path}' not found.")

# Save the presentation
presentation.save(slide_save_path)