import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt

# Create a presentation object
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw master shapes
# Rectangle
left = Inches(0.0)
top = Inches(0.0)
width = Inches(10.0)
height = Inches(0.5)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Vertical line
left = Inches(1.0)
top = Inches(0.0)
width = Inches(0.0)
height = Inches(0.5)
shape = slide.shapes.add_shape(MSO_SHAPE.LINE, left, top, width, height)
shape.line.color.rgb = RGBColor(0, 112, 192)

# Title
left = Inches(1.5)
top = Inches(0.5)
width = Inches(8.0)
height = Inches(1.0)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = 'Technical Approach: Multi-Head Attention'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Body Text
left = Inches(1.0)
top = Inches(1.5)
width = Inches(4.8)
height = Inches(4.125)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

bullet_points = [
    'Multi-head attention allows the model to attend to different parts of the input sequence simultaneously.',
    'Each head performs scaled dot-product attention.',
    'Results are concatenated and linearly transformed.'
]

for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.space_after = Pt(10)

# Image
image_path = 'output/1230_1234/images/_page_3_Figure_2.jpeg'
if os.path.exists(image_path):
    left = Inches(5.8)
    top = Inches(1.5)
    width = Inches(3.2)
    height = Inches(4.125)
    slide.shapes.add_picture(image_path, left, top, width, height)
else:
    print(f"Warning: Image file '{image_path}' not found.")

# Render LaTeX equation using matplotlib
latex_str = r'\begin{aligned} \text{MultiHead}(Q, K, V) &= \text{Concat}(\text{head}_1, ..., \text{head}_h) W^O \\ \text{where } \text{head}_{\text{i}} &= \text{Attention}(QW_i^Q, KW_i^K, VW_i^V) \end{aligned}'
fig, ax = plt.subplots(figsize=(6, 1))
ax.text(0.5, 0.5, f"${latex_str}$", fontsize=12, ha='center', va='center')
ax.axis('off')
equation_image_path = 'temp_equation.png'
fig.savefig(equation_image_path, bbox_inches='tight', pad_inches=0, dpi=300)
plt.close(fig)

# Insert LaTeX equation image
if os.path.exists(equation_image_path):
    left = Inches(1.0)
    top = Inches(5.75)
    slide.shapes.add_picture(equation_image_path, left, top)
else:
    print(f"Warning: Equation image file '{equation_image_path}' not found.")

# Save the presentation
output_path = 'output/1230_1234/result/slide_05/slide.pptx'
prs.save(output_path)
print(f"Presentation saved to '{output_path}'")