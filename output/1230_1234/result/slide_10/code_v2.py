from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a presentation object
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[5]  # 5 is typically a blank layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add "Top Header Bar"
top_header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.0), Inches(0.0), Inches(10.0), Inches(0.5)
)
top_header_bar.fill.solid()
top_header_bar.fill.fore_color.rgb = RGBColor(0, 112, 192)
top_header_bar.line.fill.background()

# Add "Vertical Line"
vertical_line = slide.shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,  # Corrected shape type
    Inches(1.0), Inches(0.0), Inches(0.0), Inches(0.5)  # Corrected dimensions for a line
)
vertical_line.line.color.rgb = RGBColor(0, 112, 192)

# Add Title
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(8.0), Inches(1.0))
title_frame = title_box.text_frame
title_frame.text = 'Questions & Discussion'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add Body Text
body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.125))
body_frame = body_box.text_frame
body_frame.word_wrap = True

# Add bullet points
bullet_points = [
    'Thank you for your attention!',
    'Questions and feedback are welcome.'
]

for point in bullet_points:
    p = body_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.alignment = PP_ALIGN.LEFT

# Ensure the first paragraph is a bullet point
body_frame.paragraphs[0].level = 0

# Save the presentation
output_path = 'output/1230_1234/result/slide_10/slide.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)