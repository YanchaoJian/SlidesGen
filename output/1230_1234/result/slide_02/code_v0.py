import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create a presentation object
    prs = Presentation()
    
    # Add a slide layout
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set slide background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add master shapes
    # Top Header Bar
    left = Inches(0.0)
    top = Inches(0.0)
    width = Inches(10.0)
    height = Inches(0.5)
    top_header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    top_header_bar.fill.solid()
    top_header_bar.fill.fore_color.rgb = RGBColor(0, 112, 192)
    
    # Vertical Line
    left = Inches(1.0)
    top = Inches(0.0)
    width = Inches(0.0)
    height = Inches(0.5)
    vertical_line = slide.shapes.add_shape(
        MSO_SHAPE.LINE, left, top, width, height
    )
    vertical_line.line.color.rgb = RGBColor(0, 112, 192)
    
    # Add title box
    left = Inches(1.5)
    top = Inches(0.5)
    width = Inches(8.0)
    height = Inches(1.0)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = 'The Problem: Limitations of Existing Methods'
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Arial'
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = RGBColor(0, 0, 0)
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add body text box
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(4.125)
    body_box = slide.shapes.add_textbox(left, top, width, height)
    body_frame = body_box.text_frame
    
    bullet_points = [
        'RNNs and CNNs face limitations in parallelization due to their sequential nature.',
        'These limitations restrict computational efficiency and scalability, especially for long sequences.',
        'A new approach is needed to overcome these constraints.'
    ]
    
    for point in bullet_points:
        p = body_frame.add_paragraph()
        p.text = point
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(169, 169, 169)
        p.space_after = Pt(12)
    
    # Save presentation
    output_path = 'output/1230_1234/result/slide_02/slide.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()