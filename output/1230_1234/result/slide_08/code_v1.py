import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Use blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw "Top Header Bar"
header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(10.0), Inches(0.5)
)
header_fill = header_bar.fill
header_fill.solid()
header_fill.fore_color.rgb = RGBColor(0, 112, 192)

# Draw "Vertical Line"
vertical_line = slide.shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER, Inches(1.0), Inches(0.0), Inches(0.0), Inches(0.5)
)
line_fill = vertical_line.line
line_fill.color.rgb = RGBColor(0, 112, 192)

# Create title box
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(8.0), Inches(1.0))
title_frame = title_box.text_frame
title_frame.text = "Ablation Study: Component Effectiveness"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Create body text box
body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(8.0), Inches(1.0))
body_frame = body_box.text_frame
body_frame.word_wrap = True

# Add bullet points
bullet_points = [
    "Ablation studies confirm the effectiveness of various model components.",
    "Multi-head attention and positional encoding are critical for performance.",
    "The model's architecture allows for efficient parallelization."
]

for point in bullet_points:
    p = body_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = Pt(21.6)  # 1.2 * 18pt

# Insert table
rows, cols = 4, 3  # Example size, adjust as needed
table_left = Inches(1.0)
table_top = Inches(2.7)  # Adjusted to ensure no overlap with text
table_width = Inches(8.0)
table_height = Inches(3.0)  # Adjusted height

table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# Set table caption
caption = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8.0), Inches(0.5))
caption_frame = caption.text_frame
caption_frame.text = "Table 3: Variations on the Transformer architecture."
caption_frame.paragraphs[0].font.size = Pt(18)
caption_frame.paragraphs[0].font.name = 'Arial'
caption_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
caption_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Set header row style
for cell in table.rows[0].cells:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
output_path = 'output/1230_1234/result/slide_08/slide.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
presentation.save(output_path)