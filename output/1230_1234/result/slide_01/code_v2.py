import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw master shapes
# Rectangle (Top Header Bar)
header_left = Inches(0.0)
header_top = Inches(0.0)
header_width = Inches(10.0)
header_height = Inches(0.5)
header_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, header_left, header_top, header_width, header_height
)
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Line (Vertical Line)
line_left = Inches(1.0)
line_top = Inches(0.0)
line_width = Inches(0.01)  # Minimal width for vertical line
line_height = Inches(0.5)
line_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height  # Use RECTANGLE for line simulation
)
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Create title box
title_left = Inches(1.5)
title_top = Inches(0.5)
title_width = Inches(8.0)
title_height = Inches(1.0)
title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_frame.text = "Background: Sequence Transduction is Transforming the World"

title_para = title_frame.paragraphs[0]
title_para.font.name = 'Arial'
title_para.font.size = Inches(28 / 72)  # Font size in points
title_para.font.color.rgb = RGBColor(0, 0, 0)
title_para.alignment = PP_ALIGN.LEFT

# Create body text box
body_left = Inches(1.0)
body_top = Inches(1.5)
body_width = Inches(8.0)
body_height = Inches(4.125)
body_box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
body_frame = body_box.text_frame

bullets = [
    "Sequence transduction is crucial for tasks like language modeling and machine translation.",
    "Traditionally, RNNs and CNNs have been the state-of-the-art methods.",
    "Attention mechanisms have been used to improve performance in these models."
]

for bullet in bullets:
    p = body_frame.add_paragraph()
    p.text = bullet
    p.font.name = 'Arial'
    p.font.size = Inches(18 / 72)  # Font size in points
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.space_before = Inches(0.1)  # Add space before text for better readability
    p.alignment = PP_ALIGN.LEFT

# Save the presentation
output_path = 'output/1230_1234/result/slide_01/slide.pptx'
if not os.path.exists(os.path.dirname(output_path)):
    os.makedirs(os.path.dirname(output_path))
presentation.save(output_path)