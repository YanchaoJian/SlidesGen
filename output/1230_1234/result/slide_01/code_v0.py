import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Use blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add "Top Header Bar"
header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.0), Inches(0.0),
    Inches(10.0), Inches(0.5)
)
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Add "Vertical Line"
vertical_line = slide.shapes.add_shape(
    MSO_SHAPE.LINE,
    Inches(1.0), Inches(0.0),
    Inches(1.0), Inches(0.5)
)
vertical_line.line.color.rgb = RGBColor(0, 112, 192)

# Add Title
title_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(0.5),
    Inches(8.0), Inches(1.0)
)
title_frame = title_box.text_frame
title_frame.text = 'Background: Sequence Transduction is Transforming the World'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add Body Text
body_box = slide.shapes.add_textbox(
    Inches(1.0), Inches(1.5),
    Inches(8.0), Inches(4.125)
)
body_frame = body_box.text_frame
body_frame.word_wrap = True

# Add bullet points
bullet_points = [
    'Sequence transduction is crucial for tasks like language modeling and machine translation.',
    'Traditionally, RNNs and CNNs have been the state-of-the-art methods.',
    'Attention mechanisms have been used to improve performance in these models.'
]

for point in bullet_points:
    p = body_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.alignment = PP_ALIGN.LEFT

# Define the output path
output_path = 'output/1230_1234/result/slide_01/slide.pptx'

# Ensure the directory exists
os.makedirs(os.path.dirname(output_path), exist_ok=True)

# Save the presentation
presentation.save(output_path)