from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a presentation object
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[5]  # 5 is a blank layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add a rectangle (Top Header Bar)
left = Inches(0.0)
top = Inches(0.0)
width = Inches(10.0)
height = Inches(0.5)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Add a vertical line
left = Inches(1.0)
top = Inches(0.0)
height = Inches(0.5)
width = Inches(0.05)  # Small width to make it visible as a vertical bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Add title box
left = Inches(1.5)
top = Inches(0.5)
width = Inches(8.0)
height = Inches(1.0)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = 'Background: Sequence Transduction is Transforming the World'
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.name = 'Arial'
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add body text box
left = Inches(1.0)
top = Inches(1.5)
width = Inches(8.0)
height = Inches(4.125)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame

# Add bullet points
bullet_points = [
    'Sequence transduction is crucial for tasks like language modeling and machine translation.',
    'Traditionally, RNNs and CNNs have been the state-of-the-art methods.',
    'Attention mechanisms have been used to improve performance in these models.'
]

for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(169, 169, 169)
    p.alignment = PP_ALIGN.LEFT

# Ensure the first paragraph is not empty
if text_frame.paragraphs[0].text == "":
    text_frame.paragraphs[0].text = bullet_points[0]

# Save the presentation
output_path = 'output/1230_1234/result/slide_01/slide.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)