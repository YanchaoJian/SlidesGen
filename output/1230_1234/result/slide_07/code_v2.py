import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Prepare presentation
prs = Presentation()

# Add slide layout
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Draw top header bar
top_header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.0),
    Inches(0.0),
    Inches(10.0),
    Inches(0.5),
)
top_header_bar.fill.solid()
top_header_bar.fill.fore_color.rgb = RGBColor(0, 112, 192)

# Draw vertical line using a shape (MSO_SHAPE.LINE is not available, using LINE instead)
slide.shapes.add_line(
    Inches(1.0),
    Inches(0.0),
    Inches(1.0),
    Inches(0.5),
).line.color.rgb = RGBColor(0, 112, 192)

# Add title
title_box = slide.shapes.add_textbox(
    Inches(1.5),
    Inches(0.5),
    Inches(8.0),
    Inches(1.0),
)
title_frame = title_box.text_frame
title_frame.text = "Key Results: BLEU Scores and Training Efficiency"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add key result text box
text_box = slide.shapes.add_textbox(
    Inches(1.0),
    Inches(1.5),
    Inches(8.0),
    Inches(2.0),
)
text_frame = text_box.text_frame
bullet_points = [
    "The Transformer achieves superior BLEU scores compared to previous models.",
    "Significantly reduced training costs and time.",
    "Demonstrates the model's efficiency and effectiveness.",
]
for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(169, 169, 169)
    if p is text_frame.paragraphs[0]:
        p.space_before = Pt(0)

# Add a table
rows, cols = 12, 5
table = slide.shapes.add_table(rows, cols, Inches(1.0), Inches(3.5), Inches(8.0), Inches(2.125)).table

# Set column headers
headers = ["Model", "BL", "EU", "Training C", "Training Cost (FLOPs)"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True

# Table data: this should be populated with actual data if available
# For demonstration, fill with placeholders
for row_idx in range(1, rows):
    for col_idx in range(cols):
        table.cell(row_idx, col_idx).text = f"Data {row_idx},{col_idx}"

# Ensure file path exists
output_path = 'output/1230_1234/result/slide_07/slide.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)

# Save presentation
prs.save(output_path)